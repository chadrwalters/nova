from pathlib import Path
import structlog
from typing import Dict, Any
from pptx import Presentation
from datetime import datetime
import aiofiles
import asyncio
from .base_processor import BaseDocumentProcessor
from ..core.errors import ProcessingError, format_error_message
from ..core.logging import get_logger

logger = get_logger(__name__)

class PowerPointProcessor(BaseDocumentProcessor):
    """Process PowerPoint documents using python-pptx"""
    
    def __init__(self, config):
        super().__init__(config)
        
        # Store config for our own use
        self.extract_images = config.document_handling.powerpoint_processing.get("extract_images", True)
        self.include_notes = config.document_handling.powerpoint_processing.get("include_notes", True)
        self.slide_separator = config.document_handling.powerpoint_processing.get("slide_separator", "---")
        self.temp_dir = config.processing.temp_dir
        
        logger.info("powerpoint_processor_initialized",
                   extract_images=self.extract_images,
                   include_notes=self.include_notes,
                   temp_dir=self.temp_dir)
        
    async def process_document(self, doc_path: Path, title: str, meta: Dict[str, Any]) -> str:
        """Process PowerPoint document and return markdown content."""
        try:
            logger.info("processing_powerpoint_document", path=str(doc_path))
            
            # Validate document exists
            if not doc_path.exists():
                return self._format_error("Document not found", title, str(doc_path))
                
            # Open PowerPoint with python-pptx
            prs = Presentation(str(doc_path))
            
            # Extract metadata
            doc_meta = self._extract_metadata(doc_path)
            doc_meta.update({
                "powerpoint_document": True,
                "processor": "pptx2md",
                "title": title,  # PowerPoint doesn't store title in core properties
                "slides": len(prs.slides)
            })
            
            # Try to extract core properties
            try:
                core_props = prs.core_properties
                doc_meta.update({
                    "author": core_props.author,
                    "category": core_props.category,
                    "comments": core_props.comments,
                    "content_status": core_props.content_status,
                    "created": core_props.created.isoformat() if core_props.created else None,
                    "modified": core_props.modified.isoformat() if core_props.modified else None,
                    "subject": core_props.subject,
                    "version": core_props.version
                })
            except Exception as e:
                logger.warning("powerpoint_metadata_extraction_failed",
                             path=str(doc_path),
                             error=str(e))
            
            # Convert PowerPoint to markdown
            markdown = []
            
            # Add metadata block
            markdown.append(f"""---
title: {doc_meta.get('title', title)}
author: {doc_meta.get('author', 'Unknown')}
date: {doc_meta.get('modified') or doc_meta.get('created') or datetime.utcnow().isoformat()}
source: {str(doc_path)}
slides: {doc_meta['slides']}
---
""")
            
            # Process each slide
            tasks = []
            for slide_num, slide in enumerate(prs.slides, 1):
                tasks.append(self._process_slide(slide, slide_num, len(prs.slides)))
            
            # Wait for all slides to be processed
            slide_contents = await asyncio.gather(*tasks)
            markdown.extend(slide_contents)
            
            return "\n".join(markdown)
            
        except Exception as e:
            logger.error("powerpoint_processing_failed",
                        path=str(doc_path),
                        error=str(e))
            return format_error_message(str(e), title, str(doc_path))
            
    async def _process_slide(self, slide, slide_num: int, total_slides: int) -> str:
        """Process a single PowerPoint slide."""
        try:
            content = []
            
            # Add slide header
            content.append(f"\n## Slide {slide_num}\n")
            
            # Process shapes on slide
            shape_tasks = []
            for shape in slide.shapes:
                shape_tasks.append(self._process_shape(shape, slide_num))
            
            # Wait for all shapes to be processed
            shape_contents = await asyncio.gather(*shape_tasks, return_exceptions=True)
            for shape_content in shape_contents:
                if isinstance(shape_content, Exception):
                    logger.warning("shape_processing_failed", error=str(shape_content))
                elif shape_content:
                    content.append(shape_content)
                    
            # Add speaker notes if enabled
            if self.include_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    content.append(f"\n> **Speaker Notes:**\n> {notes}\n")
            
            # Add slide separator if not last slide
            if slide_num < total_slides:
                content.append(f"\n{self.slide_separator}\n")
                
            return "\n".join(content)
            
        except Exception as e:
            logger.error("slide_processing_failed",
                        slide=slide_num,
                        error=str(e))
            return f"\n## Slide {slide_num}\n\n[Error processing slide: {str(e)}]\n"
            
    async def _process_shape(self, shape, slide_num: int) -> str:
        """Process a single PowerPoint shape."""
        try:
            if shape.has_text_frame:
                # Extract text
                text = shape.text.strip()
                if text:
                    return text + "\n"
                    
            elif shape.shape_type == 13:  # Picture
                if self.extract_images:
                    try:
                        # Save image to temp dir
                        img_path = Path(self.temp_dir) / f"slide_{slide_num}_img_{id(shape)}.png"
                        async with aiofiles.open(img_path, "wb") as f:
                            await f.write(shape.image.blob)
                        return f"\n![Image from slide {slide_num}]({img_path})\n"
                    except Exception as e:
                        logger.warning("image_extraction_failed",
                                     slide=slide_num,
                                     error=str(e))
            return ""
            
        except Exception as e:
            logger.warning("shape_processing_failed",
                         slide=slide_num,
                         error=str(e))
            return "" 