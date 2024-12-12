"""PowerPoint processing functionality."""

import os
from pathlib import Path
from typing import Optional, NamedTuple
import structlog
from pptx import Presentation
import aiofiles
import tempfile
import shutil

from src.core.exceptions import ProcessingError

logger = structlog.get_logger(__name__)

class ProcessedPPTX(NamedTuple):
    """Result of PowerPoint processing."""
    target_path: Optional[Path]
    title: str
    slide_count: int

class PowerPointProcessor:
    """Processes PowerPoint documents."""
    
    def __init__(self, temp_dir: Path) -> None:
        """Initialize PowerPoint processor.
        
        Args:
            temp_dir: Directory for temporary files
        """
        self.temp_dir = temp_dir
        self.logger = logger

    async def process_presentation(self, file_path: Path) -> ProcessedPPTX:
        """Process PowerPoint file to HTML.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            ProcessedPPTX containing output path and metadata
            
        Raises:
            ProcessingError: If processing fails
        """
        try:
            self.logger.info("Starting PowerPoint processing", file=str(file_path))
            
            # Create unique temp directory for this presentation
            temp_subdir = self.temp_dir / f"pptx_{file_path.stem}_{os.urandom(4).hex()}"
            temp_subdir.mkdir(parents=True, exist_ok=True)
            
            self.logger.info("Created temp directory", dir=str(temp_subdir))
            
            # Load presentation
            prs = Presentation(file_path)
            self.logger.info("Loaded presentation", slides=len(prs.slides))
            
            # Generate HTML content
            html_content = []
            html_content.append("""
                <div class="pptx-container">
                    <div class="slides-wrapper">
            """)

            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                slide_html = ['<div class="slide">']
                
                # Process shapes on slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        # Handle text content
                        slide_html.append(f'<div class="shape text-content">{shape.text}</div>')
                    if hasattr(shape, "image"):
                        # Handle images
                        image_path = temp_subdir / f"image_{i}_{hash(shape)}.png"
                        with open(image_path, 'wb') as img_file:
                            img_file.write(shape.image.blob)
                        slide_html.append(
                            f'<div class="shape image-content"><img src="{image_path}" alt="Slide {i} image"></div>'
                        )

                slide_html.append('</div>')
                html_content.extend(slide_html)

            html_content.append("</div></div>")

            # Write HTML output
            output_path = temp_subdir / f"{file_path.stem}.html"
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write('\n'.join(html_content))

            self.logger.info("PowerPoint processing complete",
                          output=str(output_path),
                          slides_processed=len(prs.slides))

            return ProcessedPPTX(
                target_path=output_path,
                title=file_path.stem,
                slide_count=len(prs.slides)
            )

        except Exception as e:
            self.logger.error("PowerPoint processing failed",
                            file=str(file_path),
                            error=str(e))
            raise ProcessingError(f"Failed to process PowerPoint file: {e}")

    async def cleanup(self) -> None:
        """Clean up temporary files."""
        try:
            if self.temp_dir.exists():
                shutil.rmtree(self.temp_dir)
        except Exception as e:
            self.logger.error("Failed to cleanup PowerPoint processor",
                            error=str(e)) 