import base64
import hashlib
import logging
import mimetypes
import os
import re
import shutil
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Dict, Optional, Union

import magic
import mammoth
import pandas as pd
import structlog
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation

from src.core.config import ProcessingConfig
from src.core.exceptions import MediaError, ProcessingError
from src.core.types import ProcessedAttachment

logger = structlog.get_logger(__name__)


class AttachmentType(Enum):
    """Types of supported attachments."""

    IMAGE = "image"
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    PDF = "pdf"
    TEXT = "text"
    UNKNOWN = "unknown"


@dataclass
class ProcessedAttachment:
    """Result of processing an attachment."""

    source_path: Path
    target_path: Optional[Path]
    html_path: Optional[Path]
    metadata: Dict
    is_valid: bool
    error: Optional[str] = None


class AttachmentProcessor:
    """Processes attachments in markdown documents."""

    SUPPORTED_EXTENSIONS = {
        # Images
        ".png": AttachmentType.IMAGE,
        ".jpg": AttachmentType.IMAGE,
        ".jpeg": AttachmentType.IMAGE,
        ".gif": AttachmentType.IMAGE,
        ".webp": AttachmentType.IMAGE,
        ".svg": AttachmentType.IMAGE,
        ".heic": AttachmentType.IMAGE,
        ".heif": AttachmentType.IMAGE,
        # Documents
        ".txt": AttachmentType.TEXT,
        ".pdf": AttachmentType.PDF,
        ".doc": AttachmentType.DOCUMENT,
        ".docx": AttachmentType.DOCUMENT,
        ".xls": AttachmentType.SPREADSHEET,
        ".xlsx": AttachmentType.SPREADSHEET,
        ".ppt": AttachmentType.PRESENTATION,
        ".pptx": AttachmentType.PRESENTATION,
    }

    # Media subdirectories for different types
    MEDIA_SUBDIRS = {
        AttachmentType.IMAGE: "images",
        AttachmentType.DOCUMENT: "documents",
        AttachmentType.SPREADSHEET: "spreadsheets",
        AttachmentType.PRESENTATION: "presentations",
        AttachmentType.PDF: "pdfs",
        AttachmentType.TEXT: "text",
        AttachmentType.UNKNOWN: "other",
    }

    def __init__(
        self,
        media_dir: Path,
        debug_dir: Optional[Path] = None,
        error_tolerance: str = "lenient",
    ):
        """Initialize the processor."""
        # Create template directory
        template_dir = media_dir.parent / "templates"
        template_dir.mkdir(parents=True, exist_ok=True)

        self.config = ProcessingConfig(
            media_dir=media_dir,
            debug_dir=debug_dir,
            error_tolerance=error_tolerance,
            template_dir=template_dir,
            relative_media_path="_media",
        )
        self.logger = structlog.get_logger(__name__)

        # Store instance variables
        self.media_dir = media_dir
        self.debug_dir = debug_dir
        self.error_tolerance = error_tolerance

        # Create directories
        self.media_dir.mkdir(parents=True, exist_ok=True)
        if self.debug_dir:
            self.debug_dir.mkdir(parents=True, exist_ok=True)
            (self.debug_dir / "attachments").mkdir(parents=True, exist_ok=True)
            (self.debug_dir / "media").mkdir(parents=True, exist_ok=True)

    def _get_media_subdir(self, file_type: AttachmentType) -> Path:
        """Get the appropriate media subdirectory for a file type."""
        subdir = self.MEDIA_SUBDIRS.get(
            file_type, self.MEDIA_SUBDIRS[AttachmentType.UNKNOWN]
        )
        media_subdir = self.media_dir / subdir
        media_subdir.mkdir(parents=True, exist_ok=True)
        return media_subdir

    def _get_relative_media_path(
        self, file_path: Path, file_type: AttachmentType
    ) -> str:
        """Get the relative path for media file references."""
        subdir = self.MEDIA_SUBDIRS.get(
            file_type, self.MEDIA_SUBDIRS[AttachmentType.UNKNOWN]
        )
        return f"_media/{subdir}/{file_path.name}"

    def process_attachment(self, file_path: Path) -> ProcessedAttachment:
        """Process an attachment file."""
        try:
            # Resolve and clean the file path
            file_path = Path(str(file_path).replace("%20", " ")).resolve()

            if not file_path.exists():
                logger.warning(f"Attachment file not found: {file_path}")
                return ProcessedAttachment(
                    source_path=file_path,
                    target_path=None,
                    html_path=None,
                    metadata={},
                    is_valid=False,
                    error="File not found",
                )

            # Determine file type
            file_type = self.SUPPORTED_EXTENSIONS.get(
                file_path.suffix.lower(), AttachmentType.UNKNOWN
            )

            # Clean up filename
            clean_name = re.sub(r"[^\w\-\.]", "_", file_path.stem)
            clean_name = re.sub(
                r"_+", "_", clean_name
            )  # Replace multiple underscores with single
            clean_name = clean_name[:50]  # Limit length

            # Generate consistent filename
            file_hash = self._compute_file_hash(file_path)
            new_filename = f"{clean_name}_{file_hash[:8]}{file_path.suffix.lower()}"

            # Get target directory and path
            target_dir = self._get_media_subdir(file_type)

            # Preserve directory structure for documents
            if file_type in [
                AttachmentType.DOCUMENT,
                AttachmentType.PDF,
                AttachmentType.PRESENTATION,
            ]:
                # Create a subdirectory using the original parent folder name
                subdir = target_dir / file_path.parent.name
                subdir.mkdir(parents=True, exist_ok=True)
                target_path = subdir / new_filename
            else:
                target_path = target_dir / new_filename

            # Copy file to target location
            target_dir.mkdir(parents=True, exist_ok=True)
            try:
                # Ensure target directory exists
                target_path.parent.mkdir(parents=True, exist_ok=True)

                # Copy the file
                shutil.copy2(file_path, target_path)
                logger.info(
                    f"Successfully copied attachment: {file_path} -> {target_path}"
                )

            except Exception as e:
                logger.error(
                    "Failed to copy file",
                    source=str(file_path),
                    target=str(target_path),
                    error=str(e),
                    exc_info=True,
                )
                raise

            # Process specific file types
            html_path = None
            if self.debug_dir and file_type in [
                AttachmentType.DOCUMENT,
                AttachmentType.SPREADSHEET,
                AttachmentType.PRESENTATION,
                AttachmentType.PDF,
            ]:
                html_path = self._convert_to_html(file_path, file_type)

            # Get relative path for media reference
            relative_path = self._get_relative_media_path(target_path, file_type)

            return ProcessedAttachment(
                source_path=file_path,
                target_path=target_path,
                html_path=html_path,
                metadata={
                    "type": file_type.value,
                    "size": file_path.stat().st_size,
                    "mime_type": mimetypes.guess_type(file_path)[0],
                    "hash": file_hash,
                    "relative_path": relative_path,
                    "original_name": file_path.name,
                },
                is_valid=True,
            )

        except Exception as e:
            logger.error(
                "Failed to process attachment",
                file=str(file_path),
                error=str(e),
                exc_info=True,
            )
            return ProcessedAttachment(
                source_path=file_path,
                target_path=None,
                html_path=None,
                metadata={},
                is_valid=False,
                error=str(e),
            )

    def _convert_to_html(
        self, file_path: Path, file_type: AttachmentType
    ) -> Optional[Path]:
        """Convert document to HTML for debug viewing."""
        try:
            # Create debug subdirectories
            html_base_dir = self.debug_dir / "html"
            html_base_dir.mkdir(parents=True, exist_ok=True)

            # Create type-specific subdirectories
            type_dirs = {
                AttachmentType.DOCUMENT: html_base_dir / "documents",
                AttachmentType.SPREADSHEET: html_base_dir / "spreadsheets",
                AttachmentType.PRESENTATION: html_base_dir / "presentations",
                AttachmentType.PDF: html_base_dir / "pdfs",
                AttachmentType.TEXT: html_base_dir / "text",
            }

            # Create subdirectories
            for dir_path in type_dirs.values():
                dir_path.mkdir(parents=True, exist_ok=True)

            # Get target directory
            html_dir = type_dirs.get(file_type)
            if not html_dir:
                return None

            # Clean up filename
            clean_name = re.sub(r"[^\w\-\.]", "_", file_path.stem)
            clean_name = re.sub(r"_+", "_", clean_name)
            clean_name = clean_name[:50]

            # Generate HTML filename
            file_hash = self._compute_file_hash(file_path)
            html_filename = f"{clean_name}_{file_hash[:8]}.html"
            html_path = html_dir / html_filename

            # Convert based on file type
            if file_type == AttachmentType.DOCUMENT:
                return self._convert_word_to_html(file_path, html_path)
            elif file_type == AttachmentType.SPREADSHEET:
                return self._convert_excel_to_html(file_path, html_path)
            elif file_type == AttachmentType.PRESENTATION:
                return self._convert_powerpoint_to_html(file_path, html_path)
            elif file_type == AttachmentType.PDF:
                return self._convert_pdf_to_html(file_path, html_path)

            return None

        except Exception as e:
            logger.error(f"Failed to convert {file_path} to HTML: {str(e)}")
            return None

    def _convert_word_to_html(self, file_path: Path, html_path: Path) -> Optional[Path]:
        """Convert Word document to HTML for debug viewing."""
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html = result.value
                self._save_debug_html(html_path, html, file_path.name)
            return html_path
        except Exception as e:
            logger.error(f"Failed to convert Word document to HTML: {str(e)}")
            return None

    def _convert_excel_to_html(
        self, file_path: Path, html_path: Path
    ) -> Optional[Path]:
        """Convert Excel file to HTML for debug viewing."""
        try:
            df = pd.read_excel(file_path)
            html = df.to_html(index=False, na_rep="", escape=False)
            self._save_debug_html(html_path, html, file_path.name)
            return html_path
        except Exception as e:
            logger.error(f"Failed to convert Excel file to HTML: {str(e)}")
            return None

    def _convert_powerpoint_to_html(
        self, file_path: Path, html_path: Path
    ) -> Optional[Path]:
        """Convert PowerPoint file to HTML for debug viewing."""
        try:
            prs = Presentation(file_path)
            html_parts = []

            for slide_number, slide in enumerate(prs.slides, 1):
                html_parts.append(f'<div class="slide"><h2>Slide {slide_number}</h2>')

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        html_parts.append(f"<p>{shape.text}</p>")

                html_parts.append("</div>")

            html = "\n".join(html_parts)
            self._save_debug_html(html_path, html, file_path.name)
            return html_path

        except Exception as e:
            logger.error(f"Failed to convert PowerPoint to HTML: {str(e)}")
            return None

    def _convert_pdf_to_html(self, file_path: Path, html_path: Path) -> Optional[Path]:
        """Convert PDF file to HTML for debug viewing."""
        try:
            # Create PDF viewer HTML
            html = f"""
            <div class="pdf-container">
                <h1>{file_path.name}</h1>
                <embed src="../pdfs/{file_path.name}" type="application/pdf" width="100%" height="600px">
                <p>If you cannot view the PDF above, <a href="../pdfs/{file_path.name}">click here to download it</a>.</p>
            </div>
            """

            # Copy PDF to debug directory
            pdf_dir = self.debug_dir / "pdfs"
            pdf_dir.mkdir(parents=True, exist_ok=True)
            shutil.copy2(file_path, pdf_dir / file_path.name)

            # Save HTML
            self._save_debug_html(html_path, html, file_path.name)
            return html_path

        except Exception as e:
            logger.error(f"Failed to convert PDF to HTML: {str(e)}")
            return None

    def _save_debug_html(self, debug_file: Path, content: str, title: str) -> None:
        """Save HTML content with proper styling."""
        html_template = """<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Helvetica, Arial, sans-serif;
            line-height: 1.6;
            padding: 2rem;
            max-width: 1200px;
            margin: 0 auto;
        }}
        .slide {{
            border: 1px solid #ddd;
            margin: 1rem 0;
            padding: 1rem;
            border-radius: 4px;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 1rem 0;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 8px;
            text-align: left;
        }}
        th {{ background-color: #f5f5f5; }}
        .pdf-container {{
            background: #f5f5f5;
            padding: 1rem;
            border-radius: 4px;
        }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    {content}
</body>
</html>"""

        debug_file.write_text(
            html_template.format(title=title, content=content), encoding="utf-8"
        )

    def _compute_file_hash(self, file_path: Path) -> str:
        """Compute MD5 hash of a file."""
        hasher = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()

    def cleanup(self) -> None:
        """Clean up temporary files."""
        try:
            shutil.rmtree(self.temp_dir)
            self.temp_dir.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            logger.error("Failed to cleanup temp files", error=str(e))

    def process_base64_attachment(
        self, base64_data: str, source_path: Path
    ) -> ProcessedAttachment:
        """Process a base64 encoded attachment."""
        try:
            # Log initial data
            logger.info(
                "Processing base64 attachment",
                source=str(source_path),
                data_length=len(base64_data),
                data_start=base64_data[:100],
            )

            # Extract image format and data
            format_match = re.match(r"data:image/(\w+);base64,(.+)", base64_data)
            if not format_match:
                logger.error(
                    "Invalid base64 image format",
                    data_start=base64_data[:100],
                    source=str(source_path),
                )
                raise ValueError("Invalid base64 image format")

            img_format, b64_data = format_match.groups()
            logger.info(
                "Found image format", format=img_format, b64_length=len(b64_data)
            )

            # Clean up base64 data - remove whitespace and newlines
            b64_data = "".join(b64_data.split())
            logger.info(
                "Cleaned base64 data",
                cleaned_length=len(b64_data),
                cleaned_start=b64_data[:100],
            )

            # Add padding if needed
            padding = len(b64_data) % 4
            if padding:
                b64_data += "=" * (4 - padding)
                logger.info(
                    "Added padding",
                    padding_added=4 - padding,
                    final_length=len(b64_data),
                )

            # Validate base64 characters
            if not re.match(r"^[A-Za-z0-9+/=]+$", b64_data):
                invalid_chars = set(re.findall(r"[^A-Za-z0-9+/=]", b64_data))
                logger.error(
                    "Invalid characters in base64 data",
                    invalid_chars=list(invalid_chars),
                    source=str(source_path),
                )
                raise ValueError("Invalid characters in base64 data")

            # Decode base64 data
            try:
                img_data = base64.b64decode(b64_data)
                logger.info(
                    "Successfully decoded base64 data", decoded_length=len(img_data)
                )
            except Exception as e:
                logger.error(
                    "Failed to decode base64 data",
                    error=str(e),
                    data_length=len(b64_data),
                    data_start=b64_data[:100],
                    source=str(source_path),
                )
                raise ValueError(f"Invalid base64 data: {str(e)}")

            # Verify image data
            if len(img_data) < 100:  # Basic size check
                logger.error(
                    "Image data too small", size=len(img_data), source=str(source_path)
                )
                raise ValueError("Image data too small")

            # Clean up source name
            clean_name = re.sub(r"[^\w\-\.]", "_", source_path.stem)
            clean_name = re.sub(r"_+", "_", clean_name)
            clean_name = clean_name[:50]

            # Generate filename
            img_hash = hashlib.md5(img_data).hexdigest()
            new_filename = f"{clean_name}_{img_hash[:8]}.{img_format.lower()}"
            logger.info("Generated filename", filename=new_filename, hash=img_hash[:8])

            # Get target directory and path
            target_dir = self._get_media_subdir(AttachmentType.IMAGE)
            target_path = target_dir / new_filename

            # Save image
            target_dir.mkdir(parents=True, exist_ok=True)
            target_path.write_bytes(img_data)

            # Verify saved file
            if not target_path.exists():
                logger.error(
                    "Failed to save image file - file does not exist",
                    target=str(target_path),
                )
                raise ValueError("Failed to save image file - file does not exist")

            if target_path.stat().st_size != len(img_data):
                logger.error(
                    "Failed to save image file - size mismatch",
                    expected=len(img_data),
                    actual=target_path.stat().st_size,
                    target=str(target_path),
                )
                raise ValueError("Failed to save image file - size mismatch")

            logger.info(
                "Successfully saved image file",
                target=str(target_path),
                size=target_path.stat().st_size,
            )

            # Get relative path for media reference
            relative_path = self._get_relative_media_path(
                target_path, AttachmentType.IMAGE
            )

            return ProcessedAttachment(
                source_path=source_path,
                target_path=target_path,
                html_path=None,
                metadata={
                    "type": AttachmentType.IMAGE.value,
                    "size": len(img_data),
                    "mime_type": f"image/{img_format.lower()}",
                    "hash": img_hash,
                    "relative_path": relative_path,
                },
                is_valid=True,
            )

        except Exception as e:
            logger.error(
                "Failed to process base64 attachment",
                error=str(e),
                source=str(source_path),
                exc_info=True,
            )
            return ProcessedAttachment(
                source_path=source_path,
                target_path=None,
                html_path=None,
                metadata={},
                is_valid=False,
                error=str(e),
            )

    def process_attachments(self, content: str, source_path: Path) -> str:
        """Process attachments in markdown content."""
        try:
            # Create attachments directory
            attachments_dir = (
                self.config.get_debug_attachments_dir()
                if self.config.debug_dir
                else Path("attachments")
            )
            attachments_dir.mkdir(parents=True, exist_ok=True)

            # Process each attachment link
            def replace_attachment(match: re.Match) -> str:
                """Replace attachment link with processed version."""
                try:
                    link_text = match.group(1)
                    file_path = match.group(2)

                    # Handle relative paths
                    if not Path(file_path).is_absolute():
                        file_path = source_path.parent / file_path

                    # Verify file exists
                    if not file_path.exists():
                        self.logger.warning(f"Attachment file not found: {file_path}")
                        return match.group(0)  # Return original link if file not found

                    # Create target directory
                    target_dir = attachments_dir / source_path.stem
                    target_dir.mkdir(parents=True, exist_ok=True)

                    # Copy file to attachments directory
                    target_file = target_dir / file_path.name
                    try:
                        import shutil

                        shutil.copy2(file_path, target_file)
                        self.logger.info(
                            f"Copied attachment: {file_path} -> {target_file}"
                        )
                    except Exception as e:
                        self.logger.error(
                            f"Failed to copy attachment {file_path}: {str(e)}"
                        )
                        return match.group(0)  # Return original link if copy fails

                    # Update link to point to copied file
                    return f"[{link_text}](../attachments/{source_path.stem}/{file_path.name})"

                except Exception as e:
                    self.logger.error(f"Failed to process attachment link: {str(e)}")
                    return match.group(0)  # Return original link if processing fails

            # Replace attachment links
            content = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", replace_attachment, content)

            return content

        except Exception as e:
            self.logger.error(f"Failed to process attachments: {str(e)}")
            return content  # Return original content if processing fails
