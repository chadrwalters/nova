from pptx import Presentation
from pathlib import Path
from typing import Dict, Any, Optional
import structlog
from datetime import datetime

from src.processors.converters.base_converter import BaseConverter
from src.core.exceptions import ConversionError

logger = structlog.get_logger(__name__)

class PowerPointConverter(BaseConverter):
    """Converts PowerPoint presentations to markdown."""
    
    async def convert(self, file_path: Path) -> Optional[str]:
        """Convert PowerPoint to markdown.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Markdown content or None if conversion fails
        """
        try:
            prs = Presentation(file_path)
            content = []
            
            # Process each slide
            for slide_number, slide in enumerate(prs.slides, 1):
                # Start slide section
                content.append(f"\n## Slide {slide_number}\n")
                
                # Get slide title if it exists
                if slide.shapes.title:
                    content.append(f"### {slide.shapes.title.text}\n")
                
                # Process shapes (text boxes, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip if this is the title we already processed
                        if shape != slide.shapes.title:
                            content.append(shape.text.strip() + "\n")
                
                # Add separator between slides
                content.append("\n---\n")
            
            return "\n".join(content)
            
        except Exception as e:
            raise ConversionError(
                f"Failed to convert PowerPoint: {str(e)}",
                details={'file': str(file_path)}
            )
    
    async def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Get PowerPoint metadata.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Dictionary of metadata
        """
        try:
            prs = Presentation(file_path)
            
            # Get core properties
            core_props = prs.core_properties
            
            return {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords or '',
                'created': core_props.created or datetime.now(),
                'modified': core_props.modified or datetime.now(),
                'last_modified_by': core_props.last_modified_by or '',
                'revision': core_props.revision or 1,
                'slide_count': len(prs.slides),
                'file_size': file_path.stat().st_size
            }
            
        except Exception as e:
            logger.warning(
                "Failed to get PowerPoint metadata",
                error=str(e),
                file=str(file_path)
            )
            return {} 