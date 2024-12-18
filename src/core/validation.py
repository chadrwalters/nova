from pathlib import Path
from typing import List, Optional, Dict, Any
import aiofiles
import logging
import re
import magic
import structlog
import os
from .errors import ValidationError, ConfigError, format_error_message
from .config import NovaConfig
import yaml
import json
import threading
from .context import set_current_frontmatter
from datetime import datetime
from .paths import resolve_path, get_nova_paths
from .errors import ErrorHandler, ProcessingError, ErrorSeverity
from .docx_processor import WordProcessor
import urllib.parse
import asyncio
from .logging import get_logger

logger = get_logger(__name__)

def _is_path_allowed(path: Path, allowed_dirs: Optional[List[Path]] = None) -> bool:
    """
    Check if path is within allowed directories.
    Args:
        path: Path to validate
        allowed_dirs: List of allowed directories. If None, uses environment variables.
    """
    try:
        path = resolve_path(path)
        
        # If no allowed dirs provided, get from environment
        if allowed_dirs is None:
            try:
                paths = get_nova_paths()
                allowed_dirs = list(paths.values())
            except ValidationError:
                return False
                
        # Resolve all allowed directories
        resolved_dirs = []
        for d in allowed_dirs:
            try:
                resolved_dirs.append(resolve_path(d))
            except Exception as e:
                logger.warning("path_resolution_failed", 
                             dir=str(d), 
                             error=str(e))
                
        # Log validation attempt
        logger.debug("path_validation",
                    check_path=str(path),
                    allowed_dirs=[str(d) for d in resolved_dirs])
        
        # Check if path is within any allowed directory
        path_str = str(path)
        for allowed_dir in resolved_dirs:
            if path_str.startswith(str(allowed_dir)):
                return True
                
        return False
        
    except Exception as e:
        logger.error("path_validation_failed", 
                    path=str(path), 
                    error=str(e))
        return False

class DocumentValidator:
    """Validates markdown documents according to Nova rules."""
    
    def __init__(self, config: Dict):
        self.config = config
        self.base_paths = self._resolve_base_paths()
        
    def _resolve_base_paths(self) -> Dict[str, Path]:
        """Resolve and validate all base paths from environment"""
        paths = {}
        for key, env_var in self.config['document_processing']['base_paths'].items():
            # Expand environment variables
            path_str = os.path.expandvars(env_var)
            path = Path(path_str)
            
            # Create directory if doesn't exist
            path.mkdir(parents=True, exist_ok=True)
            
            paths[key] = path
        return paths
        
    def validate_embedded_doc(self, doc_path: str, doc_type: str) -> Optional[Path]:
        """Validate embedded document exists and is processable"""
        # Resolve full path
        full_path = self.base_paths['input'] / doc_path
        
        # Check file exists
        if not full_path.exists():
            return None
            
        # Validate extension
        ext = full_path.suffix.lower()
        valid_extensions = self.config['document_processing']['embedded_documents'][doc_type]['extensions']
        
        if ext not in valid_extensions:
            raise ValueError(f"Invalid extension {ext} for type {doc_type}")
            
        return full_path

    def get_temp_path(self, original_path: Path) -> Path:
        """Get temporary processing path for document"""
        temp_base = self.base_paths['temp']
        return temp_base / f"{original_path.stem}_processing{original_path.suffix}"

    def get_asset_path(self, original_path: Path, asset_type: str) -> Path:
        """Generate standardized asset path"""
        date = original_path.stem.split('-')[0].strip()
        name = original_path.stem.split('-')[-1].strip()
        
        asset_name = self.config['NOVA_ASSET_NAMING_FORMAT'].format(
            date=date,
            type=asset_type,
            name=name
        )
        
        return self.base_paths['assets'] / asset_type / asset_name

    async def validate_input_files(self, files: List[Path]) -> None:
        """
        Validate all input markdown files.
        
        Args:
            files: List of file paths to validate
            
        Raises:
            ValidationError: If any validation check fails
        """
        if not files:
            self.error_handler.add_error(ProcessingError(
                message="No input files provided",
                severity=ErrorSeverity.CRITICAL,
                source="validate_input_files"
            ))
            return
            
        # Validate each file
        for file_path in files:
            await self._validate_single_file(file_path)
            
        # Validate collection as a whole
        self._validate_file_collection(files)
            
    async def _validate_single_file(self, file_path: Path) -> None:
        """Perform comprehensive validation on a single file."""
        try:
            # Basic file checks
            await self._validate_file_existence(file_path)
            await self._validate_file_extension(file_path)
            await self._validate_file_size(file_path)
            await self._validate_file_permissions(file_path)
            await self._validate_path_safety(file_path)
            
            # Content checks
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
                
            await self._validate_encoding(content)
            await self._validate_markdown_structure(content)
            await self._validate_content_safety(content)
            
        except ValidationError as e:
            logger.error("file_validation_failed",
                        file_path=str(file_path),
                        error=str(e))
            raise
        except Exception as e:
            raise ValidationError(f"Unexpected error validating {file_path}: {str(e)}")
            
    async def _validate_file_existence(self, file_path: Path) -> None:
        """Check if file exists and is a regular file."""
        if not file_path.exists():
            raise ValidationError(f"File not found: {file_path}")
        if not file_path.is_file():
            raise ValidationError(f"Not a regular file: {file_path}")
            
    async def _validate_file_extension(self, file_path: Path) -> None:
        """Validate file extension."""
        if file_path.suffix.lower() not in ['.md', '.markdown']:
            raise ValidationError(f"Invalid file extension: {file_path}")
            
    async def _validate_file_size(self, file_path: Path) -> None:
        """Check if file size is within limits."""
        if file_path.stat().st_size > self.max_file_size:
            raise ValidationError(
                f"File too large: {file_path} "
                f"({file_path.stat().st_size} bytes > {self.max_file_size} bytes)"
            )
            
    async def _validate_file_permissions(self, file_path: Path) -> None:
        """Check file permissions."""
        if not os.access(file_path, os.R_OK):
            raise ValidationError(f"File not readable: {file_path}")
            
    async def _validate_encoding(self, content: str) -> None:
        """Validate file encoding is UTF-8."""
        try:
            content.encode('utf-8')
        except UnicodeEncodeError:
            raise ValidationError("File must be UTF-8 encoded")
            
    async def _validate_markdown_structure(self, content: str) -> None:
        """Validate basic markdown structure."""
        if not content.strip():
            raise ValidationError("File is empty")
        
        # Check for balanced markdown elements but don't fail validation
        if not self._check_balanced_elements(content):
            if self.config.processing.error_tolerance == 'strict':
                raise ValidationError("Unbalanced markdown elements found")
            else:
                logger.warning("Unbalanced markdown elements found - continuing in lenient mode")
            
    async def _validate_content_safety(self, content: str) -> None:
        """Check content for potential security issues."""
        # Check for suspicious patterns but allow embedded docs
        suspicious_patterns = [
            r'<script.*?>',  # Inline JavaScript
            r'javascript:',  # JavaScript URLs
            r'data:text/html',  # Data URLs
        ]
        
        for pattern in suspicious_patterns:
            if re.search(pattern, content, re.IGNORECASE | re.MULTILINE):
                raise ValidationError("Potentially unsafe content found")
                
        # Allow HTML comments that match our embed pattern
        comment_pattern = r'<!--.*?-->'
        embed_pattern = r'<!-- \{"embed":.*?\} -->'
        
        for comment in re.finditer(comment_pattern, content, re.MULTILINE):
            if not re.match(embed_pattern, comment.group(0)):
                logger.warning("Non-embed HTML comment found")
                
    def _check_balanced_elements(self, content: str) -> bool:
        """Check if markdown elements are properly balanced."""
        # Only check code blocks - other elements can be more flexible
        backtick_count = content.count('```')
        if backtick_count % 2 != 0:
            logger.warning("Unbalanced code blocks found")
            if self.config.processing.error_tolerance == 'strict':
                return False
            
        return True
        
    def _validate_file_collection(self, files: List[Path]) -> None:
        """Validate the collection of files as a whole."""
        # Check for duplicate filenames
        filenames = [f.name for f in files]
        if len(filenames) != len(set(filenames)):
            self.error_handler.add_error(ProcessingError(
                message="Duplicate filenames found",
                severity=ErrorSeverity.ERROR,
                source="_validate_file_collection",
                details={"filenames": filenames}
            ))
            
        # Check total size of all files
        total_size = sum(f.stat().st_size for f in files)
        max_total_size = self.config.processing.max_total_size * 1024 * 1024
        if total_size > max_total_size:
            self.error_handler.add_error(ProcessingError(
                message=f"Total size of files ({total_size} bytes) exceeds limit ({max_total_size} bytes)",
                severity=ErrorSeverity.ERROR,
                source="_validate_file_collection",
                details={
                    "total_size": total_size,
                    "max_size": max_total_size,
                    "file_count": len(files)
                }
            ))
            
    async def _validate_path_safety(self, file_path: Path) -> None:
        """Validate that the file path is safe and within allowed directories."""
        try:
            if not self._is_path_allowed(file_path):
                self.error_handler.add_error(ProcessingError(
                    message=f"File {file_path.name} is outside allowed directories",
                    severity=ErrorSeverity.ERROR,
                    source="_validate_path_safety",
                    details={"file": str(file_path)}
                ))
                
            # Check for path traversal attempts
            if '..' in str(file_path):
                self.error_handler.add_error(ProcessingError(
                    message=f"Path traversal detected in {file_path.name}",
                    severity=ErrorSeverity.ERROR,
                    source="_validate_path_safety",
                    details={"file": str(file_path)}
                ))
                
        except Exception as e:
            self.error_handler.add_error(ProcessingError(
                message=f"Path validation failed for {file_path.name}: {str(e)}",
                severity=ErrorSeverity.ERROR,
                source="_validate_path_safety",
                details={"file": str(file_path), "error": str(e)}
            ))
        
    def _is_path_allowed(self, path: Path) -> bool:
        """Check if path is within allowed directories."""
        return _is_path_allowed(path, self.allowed_dirs)

    async def _validate_document_references(self, content: str) -> None:
        """Validate embedded document references."""
        doc_refs = re.findall(r'\[(.*?)\]\((.*?)\)<!--\s*({.*?})\s*-->', content)
        
        for title, path, meta_str in doc_refs:
            try:
                ref_path = Path(path)
                meta_dict = json.loads(meta_str)
                
                is_valid, error_msg, updated_meta = validate_document_reference(
                    ref_path, meta_dict, self.config
                )
                
                if not is_valid:
                    self.error_handler.add_error(ProcessingError(
                        message=error_msg,
                        severity=ErrorSeverity.WARNING,
                        source="_validate_document_references",
                        details={"ref_path": str(ref_path)}
                    ))
                    
            except json.JSONDecodeError:
                self.error_handler.add_error(ProcessingError(
                    message=f"Invalid document metadata format: {path}",
                    severity=ErrorSeverity.ERROR,
                    source="_validate_document_references",
                    details={"path": str(path)}
                ))

def load_config(config_path: Path) -> NovaConfig:
    """Load and validate configuration."""
    try:
        logger.debug("loading_config", path=str(config_path))
        with open(config_path) as f:
            config_dict = yaml.safe_load(f)
        logger.debug("config_loaded", config=config_dict)
        
        try:
            config = NovaConfig(**config_dict)
            logger.debug("config_validated", config=config.model_dump())
            return config
        except Exception as e:
            logger.error("config_validation_failed", error=str(e))
            raise ConfigError(f"Invalid configuration: {str(e)}")
    except Exception as e:
        logger.error("config_load_failed", error=str(e))
        raise ConfigError(f"Failed to load config: {str(e)}")

def handle_base64_content(content, config):
    """
    Process base64 encoded content in markdown files
    """
    # Find base64 image patterns
    base64_pattern = r'!\[.*?\]\(data:image\/[^;]+;base64,[^\)]+\)'
    
    if config.logging.filter_binary:  # Use config instead of direct env var
        # Replace with placeholder and log the size
        def replace_base64(match):
            data = match.group(0)
            size = len(data)
            return f'[BASE64 IMAGE: {size} bytes]'
            
        filtered_content = re.sub(base64_pattern, replace_base64, content)
        return filtered_content
    
    return content

async def process_embedded_docs(content: str, config: Dict, source_file: Path) -> str:
    """Process embedded documents in markdown content."""
    
    async def process_doc_match(match) -> str:
        try:
            title, path, ext, meta_str = match.groups()
            
            # URL decode the path
            decoded_path = urllib.parse.unquote(path)
            
            # Resolve full document path
            doc_path = Path(source_file.parent / decoded_path).resolve()
            
            # Log the path resolution
            logger.debug("document_path_resolution",
                        original=path,
                        decoded=decoded_path,
                        resolved=str(doc_path))
            
            # Parse metadata
            try:
                meta = json.loads(meta_str)
            except json.JSONDecodeError:
                meta = {"embed": True}
            
            # Validate document reference
            is_valid, error_msg, meta = validate_document_reference(
                doc_path, meta, config
            )
            
            if not is_valid:
                return format_error_message(error_msg, title, decoded_path)
            
            # Get appropriate processor based on file extension
            processor = None
            if ext.lower() in ['.docx', '.doc']:
                from ..processors.docx_processor import WordProcessor
                processor = WordProcessor(config)
            elif ext.lower() == '.pdf':
                from ..processors.pdf_processor import PDFProcessor
                processor = PDFProcessor(config)
            elif ext.lower() in ['.pptx', '.ppt']:
                from ..processors.pptx_processor import PowerPointProcessor
                processor = PowerPointProcessor(config)
            
            if processor:
                # Process the document
                return await processor.process_document(doc_path, title, meta)
            else:
                return format_error_message(
                    f"Unsupported document type: {ext}",
                    title,
                    decoded_path
                )
            
        except Exception as e:
            logger.error("document_processing_failed",
                        doc=str(path) if 'path' in locals() else 'unknown',
                        error=str(e))
            return format_error_message(str(e), title, decoded_path)
    
    # Find and process embedded documents
    doc_pattern = r'\[([^\]]+)\]\(([^)]+\.(docx|pdf|pptx))\)\s*<!--\s*({[^}]+})\s*-->'
    
    # Process each match asynchronously
    tasks = []
    for match in re.finditer(doc_pattern, content):
        tasks.append(process_doc_match(match))
    
    # Wait for all processing to complete
    results = await asyncio.gather(*tasks)
    
    # Replace each match with its processed content
    processed_content = content
    for match, result in zip(re.finditer(doc_pattern, content), results):
        processed_content = processed_content.replace(match.group(0), result)
    
    return processed_content

def sanitize_special_characters(content: str) -> str:
    """
    Replace problematic Unicode characters with safe alternatives
    """
    replacements = {
        '⠀': ' ',      # Braille space pattern
        ' ': ' ',      # Other special space
        '\u2028': '\n',  # Line separator
        '\u2029': '\n',  # Paragraph separator
        '\u200b': '',    # Zero-width space
        '\u200c': '',    # Zero-width non-joiner
        '\u200d': '',    # Zero-width joiner
        '\ufeff': ''     # Zero-width no-break space
    }
    
    for char, replacement in replacements.items():
        content = content.replace(char, replacement)
        
    return content

async def validate_markdown_content(content: str, config: NovaConfig, source_file: Path,
                             error_handler: ErrorHandler) -> str:
    """Validate and sanitize markdown content"""
    # Extract and validate frontmatter
    frontmatter, content = extract_frontmatter(content)
    
    # Store frontmatter in thread-local storage for metadata update
    set_current_frontmatter(frontmatter)
    
    # Sanitize special characters
    content = sanitize_special_characters(content)
    
    # Check for basic markdown structure
    if not re.search(r'^#\s+.+', content, re.MULTILINE):
        error_handler.add_error(ProcessingError(
            message="File may not have proper markdown headers",
            severity=ErrorSeverity.WARNING,
            source="validate_markdown_content",
            details={"file": str(source_file)}
        ))
    
    # Remove null bytes and control characters
    content = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', content)
    
    # Handle binary content
    content = handle_base64_content(content, config)
    
    # Process embedded documents - pass source file path and error handler
    content = await process_embedded_docs(content, config, source_file)
    
    # Reconstruct document with cleaned frontmatter if present
    if frontmatter:
        yaml_str = yaml.dump(frontmatter, default_flow_style=False, allow_unicode=True)
        content = f"---\n{yaml_str}---\n\n{content}"
    
    return content

def extract_frontmatter(content: str) -> tuple[dict, str]:
    """
    Extract and validate YAML frontmatter from markdown content.
    Returns (frontmatter_dict, remaining_content)
    """
    frontmatter_pattern = r'^---\s*\n(.*?)\n---\s*\n(.*)$'
    match = re.match(frontmatter_pattern, content, re.DOTALL)
    
    if not match:
        return {}, content
        
    try:
        # Parse YAML frontmatter with safe_load
        frontmatter_str, remaining = match.groups()
        
        # Pre-process the YAML to handle common issues
        processed_yaml = []
        for line in frontmatter_str.split('\n'):
            # Skip empty lines
            if not line.strip():
                continue
            # Fix lines with just a key
            if ':' not in line and line.strip():
                line = f"{line.strip()}: "
            # Fix lines with malformed dashes
            if ' - - ' in line:
                line = line.replace(' - - ', ': ')
            processed_yaml.append(line)
            
        frontmatter_str = '\n'.join(processed_yaml)
        frontmatter = yaml.safe_load(frontmatter_str)
        
        # Validate frontmatter structure
        if frontmatter is None:
            frontmatter = {}
        elif not isinstance(frontmatter, dict):
            logger.warning("Invalid frontmatter structure - expected dictionary")
            frontmatter = {}
            
        # Clean up frontmatter values
        cleaned = {}
        for key, value in frontmatter.items():
            # Convert None to empty string
            if value is None:
                cleaned[key] = ""
            # Convert lists/dicts to JSON strings
            elif isinstance(value, (list, dict)):
                cleaned[key] = json.dumps(value)
            else:
                cleaned[key] = str(value).strip()
                
        return cleaned, remaining.strip()
        
    except yaml.YAMLError as e:
        logger.warning(f"Invalid YAML in frontmatter: {str(e)}")
        # Return empty frontmatter but don't fail processing
        return {}, content

def validate_document_reference(ref_path: Path, meta: dict, config: NovaConfig) -> tuple[bool, str, dict]:
    """
    Validate a document reference.
    Returns (is_valid, error_message, updated_metadata)
    """
    try:
        # Validate file exists
        if not ref_path.exists():
            msg = f"Referenced document not found: {ref_path}"
            if config.processing.error_tolerance == 'strict':
                raise ValidationError(msg)
            return False, msg, meta

        # Extract and validate document type
        ext = ref_path.suffix.lower()
        doc_type = next(
            (k for k, v in config.document_handling.office_formats.items() 
             if ext in v.extensions),
            None
        )
        
        if not doc_type:
            msg = f"Unsupported document type {ext}"
            if config.processing.error_tolerance == 'strict':
                raise ValidationError(msg)
            return False, msg, meta

        # Validate metadata requirements
        required_metadata = config.document_handling.office_formats[doc_type].metadata
        
        # Extract document metadata based on type
        doc_metadata = extract_document_metadata(ref_path, doc_type)
        
        # Merge with provided metadata
        meta.update({
            "doc_type": doc_type,
            "original_metadata": doc_metadata,
            "required_metadata": required_metadata,
            "processor": get_processor_for_type(doc_type, config),
            "preview_support": config.document_handling.office_formats[doc_type].get(
                "preview_support", False
            )
        })

        # Update metadata with validation status
        meta.update({
            '_validated': True,
            '_validation_timestamp': datetime.utcnow().isoformat() + "Z",
            '_file_size': ref_path.stat().st_size,
            '_file_type': ext[1:],
            '_parent_doc': str(ref_path.parent.name),
            '_processing_status': 'pending'
        })

        return True, "", meta

    except Exception as e:
        msg = f"Validation error: {str(e)}"
        if config.processing.error_tolerance == 'strict':
            raise ValidationError(msg)
        return False, msg, meta

def extract_document_metadata(path: Path, doc_type: str) -> dict:
    """Extract metadata from document based on type"""
    metadata = {
        "filename": path.name,
        "size": path.stat().st_size,
        "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat()
    }
    
    try:
        if doc_type == "word":
            from docx import Document
            doc = Document(path)
            metadata.update({
                "author": doc.core_properties.author,
                "title": doc.core_properties.title,
                "created": doc.core_properties.created.isoformat() if doc.core_properties.created else None,
                "modified": doc.core_properties.modified.isoformat() if doc.core_properties.modified else None
            })
        elif doc_type == "pdf":
            import PyPDF2
            with open(path, 'rb') as f:
                pdf = PyPDF2.PdfReader(f)
                if pdf.metadata:
                    metadata.update(pdf.metadata)
        elif doc_type == "powerpoint":
            from pptx import Presentation
            prs = Presentation(path)
            metadata.update({
                "slide_count": len(prs.slides),
                "has_notes": any(slide.has_notes_slide for slide in prs.slides)
            })
    except Exception as e:
        logger.warning(f"Failed to extract metadata from {path}: {str(e)}")
        
    return metadata

def get_processor_for_type(doc_type: str, config: NovaConfig) -> str:
    """Get appropriate processor for document type"""
    processors = {
        "word": "markitdown",
        "pdf": "pdf2md",
        "powerpoint": "pptx2md"
    }
    return processors.get(doc_type, "unknown")

def extract_metadata_from_filename(filename: str) -> dict:
    """Extract metadata from standardized filename format YYYYMMDD - Title.md"""
    metadata = {}
    
    # Extract date if present (YYYYMMDD format)
    date_match = re.match(r'(\d{8})\s*-\s*(.+)\.md', filename)
    if date_match:
        date_str = date_match.group(1)
        # Format as YYYY-MM-DD
        metadata['date'] = f"{date_str[:4]}-{date_str[4:6]}-{date_str[6:8]}"
        metadata['title'] = date_match.group(2).strip()
    
    return metadata

def validate_frontmatter(content: str) -> dict:
    """Validate and extract YAML frontmatter"""
    frontmatter = {}
    
    # Look for YAML frontmatter between --- markers
    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    if fm_match:
        try:
            frontmatter = yaml.safe_load(fm_match.group(1))
            # Clean up any malformed fields
            if frontmatter.get('date') == '-':
                frontmatter['date'] = ''
        except yaml.YAMLError:
            logger.warning("Invalid YAML frontmatter")
    
    return frontmatter

def merge_metadata(filename_metadata: dict, frontmatter: dict, existing_meta: dict) -> dict:
    """Merge metadata from multiple sources in priority order"""
    merged = {}
    
    # Start with required fields
    merged['filename'] = existing_meta.get('filename', '')
    merged['date'] = (
        frontmatter.get('date') or 
        filename_metadata.get('date') or 
        existing_meta.get('date', '')
    )
    merged['title'] = (
        frontmatter.get('title') or
        filename_metadata.get('title') or
        existing_meta.get('title', '')
    )
    
    # Preserve existing processing history
    if '_processing_history' in existing_meta:
        merged['_processing_history'] = existing_meta['_processing_history']
    
    # Preserve TOC if it exists
    if 'toc' in existing_meta:
        merged['toc'] = existing_meta['toc']
    
    return merged