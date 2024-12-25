"""Handler for processing office documents."""

from pathlib import Path
from typing import Dict, Any, Optional, Set
import mammoth
import openpyxl
from pptx import Presentation
import tempfile

from ..base_handler import BaseHandler
from ...config.defaults import DEFAULT_CONFIG
from nova.core.utils.file_ops import FileOperationsManager

class OfficeHandler(BaseHandler):
    """Handles processing of office documents."""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the office document handler.
        
        Args:
            config: Optional configuration overrides
        """
        super().__init__(config)
        self._file_ops = FileOperationsManager()
        
        self.supported_formats = {
            # Word documents
            '.docx', '.doc',
            # PowerPoint presentations
            '.pptx', '.ppt',
            # Excel spreadsheets
            '.xlsx', '.xls',
            # PDF documents
            '.pdf'
        }
        
        # Merge default config with provided config
        self.config = {**DEFAULT_CONFIG.get('office', {}), **(config or {})}
        
        # Set up paths
        self.assets_dir = Path(self.config.get('assets_dir', ''))
        self.temp_dir = Path(self.config.get('temp_dir', ''))
        
        # Track temporary files for cleanup
        self.temp_files: Set[Path] = set()
    
    async def can_handle(self, file_path: Path) -> bool:
        """Check if file is a supported office document.
        
        Args:
            file_path: Path to the file to check
            
        Returns:
            bool: True if file is a supported office format
        """
        return (
            file_path.suffix.lower() in self.supported_formats and
            await self._file_ops.path_exists(file_path)
        )
    
    async def process(self, file_path: Path, context: Dict[str, Any]) -> Dict[str, Any]:
        """Process an office document.
        
        Args:
            file_path: Path to the office document
            context: Processing context
            
        Returns:
            Dict containing:
                - content: Extracted text content
                - metadata: Document metadata
                - assets: List of extracted assets
                - format: Document format
                - errors: List of processing errors
        """
        result = {
            'content': '',
            'metadata': {},
            'assets': [],
            'format': file_path.suffix.lower(),
            'errors': []
        }
        
        try:
            # Create required directories
            await self._file_ops.create_directory(self.assets_dir)
            await self._file_ops.create_directory(self.temp_dir)
            
            # Process based on file type
            if file_path.suffix.lower() in {'.docx', '.doc'}:
                await self._process_word(file_path, result)
            elif file_path.suffix.lower() in {'.pptx', '.ppt'}:
                await self._process_powerpoint(file_path, result)
            elif file_path.suffix.lower() in {'.xlsx', '.xls'}:
                await self._process_excel(file_path, result)
            elif file_path.suffix.lower() == '.pdf':
                await self._process_pdf(file_path, result)
            
        except Exception as e:
            result['errors'].append(str(e))
        
        return result
    
    def validate_output(self, result: Dict[str, Any]) -> bool:
        """Validate processing results.
        
        Args:
            result: Processing results to validate
            
        Returns:
            bool: True if results are valid
        """
        required_keys = {'content', 'metadata', 'assets', 'format', 'errors'}
        return (
            all(key in result for key in required_keys) and
            isinstance(result['content'], str) and
            isinstance(result['metadata'], dict) and
            isinstance(result['assets'], list) and
            isinstance(result['format'], str) and
            isinstance(result['errors'], list)
        )
    
    async def _process_word(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process Word document."""
        # Create temporary file for binary operations
        temp_file = self.temp_dir / f"temp_{file_path.stem}.docx"
        self.temp_files.add(temp_file)
        
        try:
            # Copy to temp file for processing
            await self._file_ops.copy_file(file_path, temp_file)
            
            # Convert to markdown using mammoth
            content = await self._file_ops.read_binary_file(temp_file)
            conversion = mammoth.convert_to_markdown(content)
            result['content'] = conversion.value
            result['metadata']['messages'] = conversion.messages
            
            # Extract any embedded images
            await self._extract_word_images(temp_file, result)
            
        finally:
            # Clean up temp file
            if await self._file_ops.path_exists(temp_file):
                await self._file_ops.remove_file(temp_file)
                self.temp_files.remove(temp_file)
    
    async def _process_powerpoint(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process PowerPoint presentation."""
        # Create temporary file for binary operations
        temp_file = self.temp_dir / f"temp_{file_path.stem}.pptx"
        self.temp_files.add(temp_file)
        
        try:
            # Copy to temp file for processing
            await self._file_ops.copy_file(file_path, temp_file)
            
            # Process presentation
            prs = Presentation(temp_file)
            
            # Extract slides
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"## Slide {slide_num}\n\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content += f"{shape.text}\n\n"
                
                slides.append(slide_content)
            
            result['content'] = "\n".join(slides)
            
            # Extract metadata
            result['metadata'].update({
                'slide_count': len(prs.slides),
                'title': prs.core_properties.title or ''
            })
            
        finally:
            # Clean up temp file
            if await self._file_ops.path_exists(temp_file):
                await self._file_ops.remove_file(temp_file)
                self.temp_files.remove(temp_file)
    
    async def _process_excel(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process Excel spreadsheet."""
        # Create temporary file for binary operations
        temp_file = self.temp_dir / f"temp_{file_path.stem}.xlsx"
        self.temp_files.add(temp_file)
        
        try:
            # Copy to temp file for processing
            await self._file_ops.copy_file(file_path, temp_file)
            
            # Process workbook
            wb = openpyxl.load_workbook(temp_file, data_only=True)
            
            sheets = []
            for sheet in wb.worksheets:
                sheet_content = f"### {sheet.title}\n\n"
                
                # Convert sheet to markdown table
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    rows.append(" | ".join(str(cell or '') for cell in row))
                
                if rows:
                    # Add header separator
                    rows.insert(1, "-" * len(rows[0]))
                    sheet_content += "\n".join(rows) + "\n\n"
                
                sheets.append(sheet_content)
            
            result['content'] = "\n".join(sheets)
            
            # Extract metadata
            result['metadata'].update({
                'sheet_count': len(wb.worksheets),
                'sheet_names': wb.sheetnames
            })
            
        finally:
            # Clean up temp file
            if await self._file_ops.path_exists(temp_file):
                await self._file_ops.remove_file(temp_file)
                self.temp_files.remove(temp_file)
    
    async def _process_pdf(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process PDF document."""
        # TODO: Implement PDF processing
        # This would typically use a library like PyPDF2 or pdfminer
        result['errors'].append("PDF processing not yet implemented")
    
    async def _extract_word_images(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Extract images from Word document."""
        # TODO: Implement Word image extraction
        # This would involve using python-docx to extract embedded images
        pass
    
    async def cleanup(self) -> None:
        """Clean up any resources used by the handler."""
        # Clean up any remaining temp files
        for temp_file in self.temp_files.copy():
            try:
                if await self._file_ops.path_exists(temp_file):
                    await self._file_ops.remove_file(temp_file)
                self.temp_files.remove(temp_file)
            except Exception as e:
                self.logger.error(f"Failed to remove temp file {temp_file}: {e}") 