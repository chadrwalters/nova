"""Handler for processing office documents."""

import os
from pathlib import Path
from typing import Dict, Any, Optional, Set
import mammoth
import openpyxl
from pptx import Presentation
import aiofiles

from ..base_handler import BaseHandler
from ...config.defaults import DEFAULT_CONFIG

class OfficeHandler(BaseHandler):
    """Handles processing of office documents."""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the office document handler.
        
        Args:
            config: Optional configuration overrides
        """
        super().__init__(config)
        self.supported_formats = {
            # Word documents
            '.docx', '.doc',
            # PowerPoint presentations
            '.pptx', '.ppt',
            # Excel spreadsheets
            '.xlsx', '.xls',
            # PDF documents
            '.pdf'
        }
        
        # Merge default config with provided config
        self.config = {**DEFAULT_CONFIG.get('office', {}), **(config or {})}
        
        # Set up paths
        self.assets_dir = Path(os.getenv('NOVA_OFFICE_ASSETS_DIR', ''))
        self.temp_dir = Path(os.getenv('NOVA_OFFICE_TEMP_DIR', ''))
    
    def can_handle(self, file_path: Path) -> bool:
        """Check if file is a supported office document.
        
        Args:
            file_path: Path to the file to check
            
        Returns:
            bool: True if file is a supported office format
        """
        return file_path.suffix.lower() in self.supported_formats
    
    async def process(self, file_path: Path, context: Dict[str, Any]) -> Dict[str, Any]:
        """Process an office document.
        
        Args:
            file_path: Path to the office document
            context: Processing context
            
        Returns:
            Dict containing:
                - content: Extracted text content
                - metadata: Document metadata
                - assets: List of extracted assets
                - format: Document format
                - errors: List of processing errors
        """
        result = {
            'content': '',
            'metadata': {},
            'assets': [],
            'format': file_path.suffix.lower(),
            'errors': []
        }
        
        try:
            if file_path.suffix.lower() in {'.docx', '.doc'}:
                await self._process_word(file_path, result)
            elif file_path.suffix.lower() in {'.pptx', '.ppt'}:
                await self._process_powerpoint(file_path, result)
            elif file_path.suffix.lower() in {'.xlsx', '.xls'}:
                await self._process_excel(file_path, result)
            elif file_path.suffix.lower() == '.pdf':
                await self._process_pdf(file_path, result)
            
        except Exception as e:
            result['errors'].append(str(e))
        
        return result
    
    def validate_output(self, result: Dict[str, Any]) -> bool:
        """Validate processing results.
        
        Args:
            result: Processing results to validate
            
        Returns:
            bool: True if results are valid
        """
        required_keys = {'content', 'metadata', 'assets', 'format', 'errors'}
        return (
            all(key in result for key in required_keys) and
            isinstance(result['content'], str) and
            isinstance(result['metadata'], dict) and
            isinstance(result['assets'], list) and
            isinstance(result['format'], str) and
            isinstance(result['errors'], list)
        )
    
    async def _process_word(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process Word document."""
        # Convert to markdown using mammoth
        with open(file_path, 'rb') as docx_file:
            conversion = mammoth.convert_to_markdown(docx_file)
            result['content'] = conversion.value
            result['metadata']['messages'] = conversion.messages
        
        # Extract any embedded images
        await self._extract_word_images(file_path, result)
    
    async def _process_powerpoint(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process PowerPoint presentation."""
        prs = Presentation(file_path)
        
        # Extract slides
        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = f"## Slide {slide_num}\n\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content += f"{shape.text}\n\n"
            
            slides.append(slide_content)
        
        result['content'] = "\n".join(slides)
        
        # Extract metadata
        result['metadata'].update({
            'slide_count': len(prs.slides),
            'title': prs.core_properties.title or ''
        })
    
    async def _process_excel(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process Excel spreadsheet."""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        sheets = []
        for sheet in wb.worksheets:
            sheet_content = f"### {sheet.title}\n\n"
            
            # Convert sheet to markdown table
            rows = []
            for row in sheet.iter_rows(values_only=True):
                rows.append(" | ".join(str(cell or '') for cell in row))
            
            if rows:
                # Add header separator
                rows.insert(1, "-" * len(rows[0]))
                sheet_content += "\n".join(rows) + "\n\n"
            
            sheets.append(sheet_content)
        
        result['content'] = "\n".join(sheets)
        
        # Extract metadata
        result['metadata'].update({
            'sheet_count': len(wb.worksheets),
            'sheet_names': wb.sheetnames
        })
    
    async def _process_pdf(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Process PDF document."""
        # TODO: Implement PDF processing
        # This would typically use a library like PyPDF2 or pdfminer
        result['errors'].append("PDF processing not yet implemented")
    
    async def _extract_word_images(self, file_path: Path, result: Dict[str, Any]) -> None:
        """Extract images from Word document."""
        # TODO: Implement Word image extraction
        # This would involve using python-docx to extract embedded images
        pass 