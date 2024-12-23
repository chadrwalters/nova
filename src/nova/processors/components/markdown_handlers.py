"""Markdown handling components for Nova processors."""

from pathlib import Path
from typing import Dict, Any, Optional, List
import re
from datetime import datetime
from urllib.parse import unquote, quote
import os
import shutil
from PIL import Image

from . import MarkdownComponent
from ...core.config import NovaConfig, ProcessorConfig
from ...core.errors import ProcessingError
from ...core.logging import get_logger
from ..image_processor import ImageProcessor

try:
    import pillow_heif
    HEIF_SUPPORT = True
    pillow_heif.register_heif_opener()
except ImportError:
    HEIF_SUPPORT = False

class MarkitdownHandler(MarkdownComponent):
    """Handles markdown processing using markitdown."""
    
    def __init__(self, processor_config: ProcessorConfig, nova_config: NovaConfig, image_processor: Optional[ImageProcessor] = None):
        """Initialize handler.
        
        Args:
            processor_config: Processor-specific configuration
            nova_config: Global Nova configuration
            image_processor: Optional image processor instance
        """
        super().__init__(processor_config, nova_config)
        self.logger = get_logger(self.__class__.__name__)
        self.image_processor = image_processor
        
        # Add component-specific stats
        self.stats.update({
            'files_processed': 0,
            'images_processed': 0,
            'links_updated': 0
        })
    
    def process_markdown(self, content: str, source_path: Path) -> str:
        """Process markdown content.
        
        Args:
            content: Markdown content to process
            source_path: Path to source file
            
        Returns:
            Processed markdown content
        """
        try:
            self.logger.info(f"Processing markdown file: {source_path}")
            
            # Process images in the content
            self.logger.info("Processing images in markdown content")
            
            # Regular expression to find image references
            # Handle both ![alt](path) and ![](path) formats, with optional whitespace
            image_pattern = r'^\*\s*(?:JPG|HEIC)(?::\s*|\s+)!\[(.*?)\]\(([^)]+)\)(?:\s*<!--\s*(\{.*?\})\s*-->)?$'
            
            # Process all images line by line
            lines = []
            for line in content.splitlines():
                # Debug logging
                if 'JPG' in line or 'HEIC' in line:
                    self.logger.info(f"Processing line: {line}")
                    match = re.match(image_pattern, line)
                    if match:
                        self.logger.info("Line matched pattern")
                        lines.append(self._process_image_match(match, source_path))
                    else:
                        self.logger.info("Line did not match pattern")
                        lines.append(line)
                else:
                    # Preserve original markdown line
                    lines.append(line)
            
            # Update stats
            self.stats['files_processed'] += 1
            
            # Return processed content while preserving markdown format
            return '\n'.join(lines)
            
        except Exception as e:
            self.logger.error(f"Failed to process {source_path}: {e}")
            raise ProcessingError(f"Failed to process {source_path}: {e}") from e
    
    def convert_to_markdown(self, input_path: Path, output_path: Optional[Path] = None) -> str:
        """Convert a file to markdown format.
        
        Args:
            input_path: Path to input file
            output_path: Optional path to save markdown output
            
        Returns:
            Converted markdown content
        """
        try:
            self.logger.info(f"Converting file to markdown: {input_path}")
            
            # Get file type from suffix
            suffix = input_path.suffix.lower()
            
            # Convert based on file type
            if suffix in ['.docx', '.doc']:
                self.logger.info(f"Converting Word document: {input_path}")
                content = self._convert_docx(input_path, output_path)
            elif suffix in ['.pptx', '.ppt']:
                self.logger.info(f"Converting PowerPoint presentation: {input_path}")
                content = self._convert_pptx(input_path, output_path)
            elif suffix in ['.xlsx', '.xls']:
                self.logger.info(f"Converting Excel workbook: {input_path}")
                content = self._convert_xlsx(input_path, output_path)
            elif suffix == '.pdf':
                self.logger.info(f"Converting PDF document: {input_path}")
                content = self._convert_pdf(input_path, output_path)
            elif suffix == '.csv':
                self.logger.info(f"Converting CSV file: {input_path}")
                content = self._convert_csv(input_path, output_path)
            elif suffix in ['.txt', '.json', '.env', '.html']:
                self.logger.info(f"Converting text file: {input_path}")
                # For text files, read directly
                with open(input_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif suffix in ['.jpg', '.jpeg', '.png', '.gif', '.heic']:
                self.logger.info(f"Creating markdown reference for image: {input_path}")
                # For images, create a markdown image reference
                content = f"![{input_path.stem}]({input_path.name})"
            else:
                error = f"Unsupported file format '{suffix}' for file: {input_path}"
                self.logger.error(error)
                raise ProcessingError(error)
            
            # Save to output file if specified
            if output_path:
                self.logger.info(f"Writing markdown output to: {output_path}")
                output_path.parent.mkdir(parents=True, exist_ok=True)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(content)
            
            self.logger.info(f"Successfully converted {input_path} to markdown")
            return content
            
        except Exception as e:
            error = f"Failed to convert {input_path} (format: {suffix}): {str(e)}"
            self.logger.error(error)
            return f"# Error Converting {input_path.name}\n\nFile type: {suffix}\nError: {str(e)}\n"
    
    def _convert_docx(self, input_path: Path, output_path: Path) -> str:
        """Convert DOCX to markdown."""
        from docx import Document
        
        self.logger.info(f"Opening Word document: {input_path}")
        doc = Document(str(input_path))
        content = []
        
        self.logger.info(f"Extracting paragraphs from: {input_path}")
        for para in doc.paragraphs:
            content.append(para.text)
        
        self.logger.info(f"Successfully extracted {len(content)} paragraphs from: {input_path}")
        return '\n\n'.join(content)
    
    def _convert_pptx(self, input_path: Path, output_path: Path) -> str:
        """Convert PPTX to markdown."""
        from pptx import Presentation
        
        self.logger.info(f"Opening PowerPoint presentation: {input_path}")
        prs = Presentation(str(input_path))
        content = []
        
        self.logger.info(f"Processing slides from: {input_path}")
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            self.logger.debug(f"Processing slide {i} from: {input_path}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            content.append('\n'.join(slide_content))
        
        self.logger.info(f"Successfully processed {len(prs.slides)} slides from: {input_path}")
        return '\n\n'.join(content)
    
    def _convert_xlsx(self, input_path: Path, output_path: Path) -> str:
        """Convert XLSX to markdown."""
        from openpyxl import load_workbook
        from datetime import datetime
        
        def format_cell_value(value):
            """Format cell value for markdown."""
            if value is None:
                return ''
            elif isinstance(value, datetime):
                return value.strftime('%Y-%m-%d')  # Format dates as YYYY-MM-DD
            else:
                return str(value).strip()
        
        self.logger.info(f"Opening Excel workbook: {input_path}")
        wb = load_workbook(str(input_path))
        content = []
        
        self.logger.info(f"Processing {len(wb.sheetnames)} sheets from: {input_path}")
        for sheet in wb:
            self.logger.debug(f"Processing sheet '{sheet.title}' from: {input_path}")
            # Add sheet name as header
            content.append(f"# {sheet.title}\n")
            
            # Process each sheet
            row_idx = 1
            tables_found = 0
            while row_idx <= sheet.max_row:
                # Find the start of a table (non-empty row)
                if any(cell.value for cell in sheet[row_idx]):
                    tables_found += 1
                    self.logger.debug(f"Found table {tables_found} in sheet '{sheet.title}' at row {row_idx}")
                    table_content = []
                    header_cells = list(sheet[row_idx])
                    
                    # Look ahead for merged header rows
                    next_row_idx = row_idx + 1
                    while next_row_idx <= sheet.max_row:
                        next_row = list(sheet[next_row_idx])
                        if not any(cell.value for cell in next_row):
                            break
                        if all(isinstance(cell.value, (str, type(None))) for cell in next_row):
                            header_cells.extend(next_row)
                            next_row_idx += 1
                        else:
                            break
                    
                    # Build header from merged rows
                    headers = []
                    for cell in header_cells:
                        value = format_cell_value(cell.value)
                        if value:
                            headers.append(value)
                    
                    # Skip if no valid headers
                    if not headers:
                        row_idx = next_row_idx
                        continue
                    
                    # Add header row
                    table_content.append(' | '.join(headers))
                    
                    # Add separator row
                    table_content.append('|'.join(['---' for _ in headers]))
                    
                    # Process data rows
                    data_start_row = next_row_idx
                    rows_processed = 0
                    while data_start_row <= sheet.max_row:
                        row = list(sheet[data_start_row])
                        row_content = []
                        is_empty_row = True
                        
                        # Process each cell in the row
                        for cell in row[:len(headers)]:  # Only process cells up to header count
                            value = format_cell_value(cell.value)
                            row_content.append(value)
                            if value:
                                is_empty_row = False
                        
                        # Skip empty rows
                        if is_empty_row:
                            data_start_row += 1
                            continue
                        
                        # Pad row if needed
                        while len(row_content) < len(headers):
                            row_content.append('')
                        
                        # Add data row
                        table_content.append(' | '.join(row_content))
                        data_start_row += 1
                        rows_processed += 1
                    
                    self.logger.debug(f"Processed {rows_processed} rows in table {tables_found}")
                    # Add table to content with spacing
                    content.extend(['', *table_content, ''])
                    row_idx = data_start_row
                else:
                    row_idx += 1
            
            self.logger.debug(f"Found {tables_found} tables in sheet '{sheet.title}'")
            # Add blank line between sheets
            content.append('')
        
        self.logger.info(f"Successfully processed workbook with {len(wb.sheetnames)} sheets from: {input_path}")
        return '\n'.join(content)
    
    def _convert_pdf(self, input_path: Path, output_path: Path) -> str:
        """Convert PDF to markdown."""
        import warnings
        from pypdf import PdfReader
        
        self.logger.info(f"Opening PDF document: {input_path}")
        # Suppress PyPDF deprecation warning about ARC4
        with warnings.catch_warnings():
            warnings.filterwarnings(
                'ignore',
                message='ARC4 has been moved to cryptography.hazmat.decrepit.*',
                category=DeprecationWarning
            )
            reader = PdfReader(str(input_path))
            content = []
            
            self.logger.info(f"Processing {len(reader.pages)} pages from: {input_path}")
            for i, page in enumerate(reader.pages, 1):
                self.logger.debug(f"Extracting text from page {i} of: {input_path}")
                content.append(page.extract_text())
            
            self.logger.info(f"Successfully extracted text from {len(reader.pages)} pages of: {input_path}")
            return '\n\n'.join(content)
    
    def _convert_csv(self, input_path: Path, output_path: Path) -> str:
        """Convert CSV to markdown."""
        import csv
        import chardet
        
        # Detect encoding
        self.logger.info(f"Detecting encoding for CSV file: {input_path}")
        with open(input_path, 'rb') as f:
            raw = f.read()
            result = chardet.detect(raw)
            encoding = result['encoding']
        self.logger.debug(f"Detected encoding {encoding} for: {input_path}")
        
        # Read CSV with detected encoding
        content = []
        self.logger.info(f"Reading CSV file with encoding {encoding}: {input_path}")
        with open(input_path, 'r', encoding=encoding) as f:
            reader = csv.reader(f)
            rows_processed = 0
            for row in reader:
                content.append(' | '.join(row))
                rows_processed += 1
        
        self.logger.info(f"Successfully processed {rows_processed} rows from: {input_path}")
        return '\n'.join(content)
    
    def _validate_image_file(self, image_path: Path) -> tuple[bool, str]:
        """Validate an image file.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        self.logger.debug(f"Validating image file: {image_path}")
        
        if not image_path.exists():
            self.logger.debug(f"Image file not found: {image_path}")
            return False, f"Image file not found: {image_path}"
            
        if not image_path.is_file():
            self.logger.debug(f"Path exists but is not a file: {image_path}")
            return False, f"Path exists but is not a file: {image_path}"
            
        # Check if it's a supported image format
        suffix = image_path.suffix.lower()
        supported_formats = ['.jpg', '.jpeg', '.png', '.gif', '.heic', '.webp']
        if suffix not in supported_formats:
            self.logger.debug(f"Unsupported image format '{suffix}' for file: {image_path}")
            return False, f"Unsupported image format '{suffix}' for file: {image_path}"
            
        # Check if file is readable
        try:
            image_path.open('rb').close()
            self.logger.debug(f"Successfully validated image file: {image_path}")
            return True, ""
        except Exception as e:
            self.logger.debug(f"File is not readable: {image_path} ({str(e)})")
            return False, f"File is not readable: {image_path} ({str(e)})"
    
    def _convert_heic_to_jpeg(self, input_path: Path, output_path: Path) -> Optional[Path]:
        """Convert HEIC image to JPEG.
        
        Args:
            input_path: Path to HEIC image
            output_path: Path to save JPEG image
            
        Returns:
            Path to converted JPEG file or None if conversion failed
        """
        try:
            self.logger.debug(f"Converting HEIC to JPEG: {input_path} -> {output_path}")
            
            if not HEIF_SUPPORT:
                self.logger.warning("HEIC conversion not available - pillow-heif not installed")
                return None
                
            # Open HEIC image using pillow-heif
            heif_file = pillow_heif.read_heif(str(input_path))
            image = Image.frombytes(
                heif_file.mode,
                heif_file.size,
                heif_file.data,
                "raw",
                heif_file.mode,
                heif_file.stride,
            )
            
            # Save as JPEG
            output_path.parent.mkdir(parents=True, exist_ok=True)
            image.save(output_path, 'JPEG', quality=85)
            self.logger.debug(f"Successfully converted HEIC to JPEG: {output_path}")
            
            return output_path
            
        except Exception as e:
            self.logger.error(f"Failed to convert HEIC to JPEG: {str(e)}")
            self.logger.debug(f"Conversion error details: {type(e).__name__}: {str(e)}")
            return None
    
    def _process_image_match(self, match: re.Match, source_path: Path) -> str:
        """Process a single image match and return updated markdown.
        
        Args:
            match: Regular expression match object containing image information
            source_path: Path to the source markdown file
            
        Returns:
            Updated markdown with processed image reference
        """
        try:
            alt_text, image_path, metadata_json = match.groups()
            self.logger.info(f"Processing image: {image_path}")
            self.logger.debug(f"Source markdown file: {source_path}")
            self.logger.debug(f"Alt text: {alt_text}")
            self.logger.debug(f"Metadata: {metadata_json}")
            
            # URL decode the image path
            try:
                decoded_path = unquote(image_path)
                if decoded_path != image_path:
                    self.logger.debug(f"Decoded image path from {image_path} to {decoded_path}")
                image_path = decoded_path
            except Exception as e:
                self.logger.warning(f"URL decoding failed for {image_path}: {str(e)}")
                self.logger.debug(f"URL decoding error details: {type(e).__name__}: {str(e)}")
            
            # Resolve image path relative to markdown file
            full_image_path = source_path.parent / Path(image_path)
            self.logger.debug(f"Resolved full image path: {full_image_path}")
            self.logger.debug(f"Source parent directory: {source_path.parent}")
            
            # Check multiple possible locations for the image
            possible_paths = [
                full_image_path,  # Original path
                source_path.parent / Path(image_path).name,  # Just filename in same dir
                Path(os.getenv('NOVA_ORIGINAL_IMAGES_DIR', '')) / Path(image_path).name,  # Original images dir
                Path(os.getenv('NOVA_PROCESSED_IMAGES_DIR', '')) / Path(image_path).name,  # Processed images dir
            ]
            
            self.logger.debug("Checking possible image locations:")
            for i, path in enumerate(possible_paths, 1):
                self.logger.debug(f"  {i}. {path}")
            
            # Try to find and validate the image in any of the possible locations
            image_found = False
            found_path = None
            validation_errors = []
            
            for path in possible_paths:
                self.logger.debug(f"Checking path: {path}")
                is_valid, error = self._validate_image_file(path)
                if is_valid:
                    image_found = True
                    found_path = path
                    self.logger.debug(f"Found valid image at: {path}")
                    break
                else:
                    validation_errors.append(f"  - {path}: {error}")
                    self.logger.debug(f"Validation failed: {error}")
            
            if not image_found:
                # Log all validation errors for debugging
                self.logger.warning(f"No valid image found for: {image_path}")
                self.logger.debug("Validation results:")
                for error in validation_errors:
                    self.logger.debug(error)
                return match.group(0)
            
            # Process image
            output_dir = source_path.parent / Path(image_path).parent
            self.logger.debug(f"Creating output directory: {output_dir}")
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Handle HEIC conversion if needed
            if found_path.suffix.lower() == '.heic':
                self.logger.debug("HEIC image detected, converting to JPEG")
                jpeg_path = output_dir / f"{found_path.stem}.jpg"
                converted_path = self._convert_heic_to_jpeg(found_path, jpeg_path)
                if converted_path:
                    found_path = converted_path
                    self.logger.debug(f"Using converted JPEG: {found_path}")
                else:
                    self.logger.warning("HEIC conversion failed, using original file")
            
            # Copy image to output directory if needed
            output_path = output_dir / found_path.name
            self.logger.debug(f"Target output path: {output_path}")
            
            if found_path != output_path:
                try:
                    self.logger.debug(f"Copying image from {found_path} to {output_path}")
                    shutil.copy2(found_path, output_path)
                    self.logger.debug(f"Successfully copied image")
                except Exception as e:
                    self.logger.error(f"Failed to copy image: {str(e)}")
                    self.logger.debug(f"Copy error details: {type(e).__name__}: {str(e)}")
                    return match.group(0)
            else:
                self.logger.debug("Image already in correct location, no copy needed")
            
            # Get the processed image path relative to the markdown file
            rel_path = output_path.relative_to(source_path.parent)
            self.logger.debug(f"Relative path from markdown: {rel_path}")
            
            # URL encode the path properly
            try:
                # Split path into parts and encode each part separately
                encoded_parts = [quote(part) for part in str(rel_path).split('/')]
                rel_path_str = '/'.join(encoded_parts)
                self.logger.debug(f"Encoded path: {rel_path_str}")
            except Exception as e:
                self.logger.warning(f"URL encoding failed for {rel_path}: {str(e)}")
                self.logger.debug(f"URL encoding error details: {type(e).__name__}: {str(e)}")
                rel_path_str = str(rel_path)
            
            # Use original alt text if not empty, otherwise use filename
            alt = alt_text if alt_text else Path(image_path).stem
            prefix = match.group(0).split('!')[0]  # Get the "* JPG:" or "* HEIC:" part
            
            # Update stats
            self.stats['images_processed'] += 1
            self.logger.debug(f"Updated stats - images processed: {self.stats['images_processed']}")
            
            result = f"{prefix}![{alt}]({rel_path_str})"
            self.logger.debug(f"Generated markdown: {result}")
            return result
            
        except Exception as e:
            error = f"Failed to process image match: {str(e)}"
            self.logger.warning(error)
            self.logger.debug(f"Error details: {type(e).__name__}: {str(e)}")
            return match.group(0)

class ConsolidationHandler(MarkdownComponent):
    """Handles markdown consolidation."""
    
    def __init__(self, processor_config: ProcessorConfig, nova_config: NovaConfig):
        """Initialize handler.
        
        Args:
            processor_config: Processor-specific configuration
            nova_config: Global Nova configuration
        """
        super().__init__(processor_config, nova_config)
        self.logger = get_logger(self.__class__.__name__)
        
        # Add component-specific stats
        self.stats.update({
            'files_consolidated': 0,
            'total_files': 0,
            'total_size': 0
        })
    
    def process_markdown(self, content: str, source_path: Path) -> str:
        """Process markdown content.
        
        Args:
            content: Markdown content to process
            source_path: Path to source file
            
        Returns:
            Processed markdown content
        """
        # This handler doesn't modify individual markdown files
        return content
    
    def consolidate_markdown(self, input_files: List[Path], output_path: Path) -> Path:
        """Consolidate markdown files.
        
        Args:
            input_files: List of input files to consolidate
            output_path: Path to output file
            
        Returns:
            Path to consolidated file
        """
        try:
            # Sort files by date in filename (YYYYMMDD format)
            date_pattern = r'(\d{8})'
            sorted_files = sorted(
                input_files,
                key=lambda x: re.search(date_pattern, x.name).group(1) if re.search(date_pattern, x.name) else '00000000'
            )
            
            # Read and process each file
            consolidated_content = []
            for file_path in sorted_files:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    
                    # Add file header
                    header = f"\n\n## {file_path.stem}\n\n"
                    consolidated_content.append(header)
                    
                    # Add processed content
                    consolidated_content.append(content)
                    
                    # Update stats
                    self.stats['files_consolidated'] += 1
                    self.stats['total_size'] += len(content)
            
            # Write consolidated content
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(consolidated_content))
            
            # Update total files stat
            self.stats['total_files'] = len(input_files)
            
            return output_path
            
        except Exception as e:
            raise ProcessingError(f"Failed to consolidate markdown: {e}") from e