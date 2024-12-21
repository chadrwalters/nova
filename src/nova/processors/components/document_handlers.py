"""Document handling components for Nova processors."""

from pathlib import Path
import json
from typing import Dict, Any, Optional, List
from markdown_it import MarkdownIt
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from . import DocumentComponent
from ...core.errors import DocumentProcessingError
from ...core.config import NovaConfig
from ...core.logging import get_logger

class OfficeDocumentHandler(DocumentComponent):
    """Handles office document processing."""
    
    def __init__(self, config: NovaConfig):
        """Initialize handler."""
        super().__init__(config)
        self.logger = get_logger(self.__class__.__name__)
        
        # Initialize markdown parser for text conversion
        self.md = MarkdownIt('commonmark', {
            'typographer': True,
            'html': True,
            'linkify': True
        })
        
        # Add component-specific stats
        self.stats.update({
            'text_extracted': 0,
            'images_extracted': 0,
            'metadata_preserved': 0
        })
    
    def process_document(self, input_path: Path, output_path: Path) -> Dict[str, Any]:
        """Process an office document."""
        try:
            # Extract text content
            text_content = self._extract_text(input_path)
            
            # Extract images if enabled
            images = []
            if self.config.processors["office"]["image_extraction"]["enabled"]:
                images = self._extract_images(input_path)
            
            # Convert text to markdown
            markdown_content = self._convert_to_markdown(text_content)
            
            # Write output file
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            # Save metadata
            metadata = {
                'original_path': str(input_path),
                'text_content': text_content,
                'images': [str(img) for img in images],
                'format': input_path.suffix.lower()
            }
            
            return metadata
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to process {input_path}: {e}") from e
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text from document."""
        suffix = file_path.suffix.lower()
        
        if suffix in ['.docx', '.doc']:
            return self._extract_docx_text(file_path)
        elif suffix in ['.pptx', '.ppt']:
            return self._extract_pptx_text(file_path)
        elif suffix in ['.xlsx', '.xls']:
            return self._extract_xlsx_text(file_path)
        elif suffix == '.pdf':
            return self._extract_pdf_text(file_path)
        else:
            raise DocumentProcessingError(f"Unsupported document format: {suffix}")
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        doc = Document(file_path)
        return '\n\n'.join(p.text for p in doc.paragraphs)
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text.append('\n'.join(slide_text))
        return '\n\n'.join(text)
    
    def _extract_xlsx_text(self, file_path: Path) -> str:
        """Extract text from XLSX file."""
        wb = load_workbook(file_path)
        text = []
        for sheet in wb:
            sheet_text = []
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value:
                        row_text.append(str(cell.value))
                if row_text:
                    sheet_text.append(' | '.join(row_text))
            text.append('\n'.join(sheet_text))
        return '\n\n'.join(text)
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return '\n\n'.join(text)
    
    def _extract_images(self, file_path: Path) -> List[Path]:
        """Extract images from document."""
        # TODO: Implement image extraction
        return []
    
    def _convert_to_markdown(self, text: str) -> str:
        """Convert text to markdown format."""
        # Basic conversion - could be enhanced
        lines = text.split('\n')
        markdown = []
        
        for line in lines:
            line = line.strip()
            if not line:
                markdown.append('')
            elif line.endswith(':'):
                markdown.append(f'### {line}')
            else:
                markdown.append(line)
        
        return '\n'.join(markdown)