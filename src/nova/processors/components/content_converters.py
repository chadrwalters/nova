"""Content conversion handlers for different file types."""

import os
import re
import json
import csv
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple, Union
from bs4 import BeautifulSoup
import markdown
import pandas as pd
from docx import Document
import fitz  # PyMuPDF
from PIL import Image
import yaml
import docx
import openpyxl
import pptx
import datetime

from ...core.errors import ProcessingError
from ...core.logging import get_logger

logger = get_logger(__name__)

class BaseContentConverter:
    """Base class for content converters."""
    
    def __init__(self):
        """Initialize the converter."""
        self.logger = get_logger(self.__class__.__name__)
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert file content to markdown.
        
        Args:
            file_path: Path to the file to convert
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        raise NotImplementedError("Subclasses must implement convert()")

class HTMLConverter(BaseContentConverter):
    """Converts HTML content to markdown."""
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert HTML to markdown.
        
        Args:
            file_path: Path to HTML file
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read HTML content
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract title
            title = soup.title.string if soup.title else os.path.basename(file_path)
            
            # Extract metadata
            metadata = {
                'title': title,
                'original_format': 'html',
                'tables': len(soup.find_all('table')),
                'images': len(soup.find_all('img')),
                'links': len(soup.find_all('a'))
            }
            
            # Clean up HTML
            # Remove script and style elements
            for element in soup(['script', 'style']):
                element.decompose()
            
            # Convert tables to markdown format
            for table in soup.find_all('table'):
                markdown_table = ['|']
                
                # Get headers
                headers = table.find_all('th')
                if not headers:
                    headers = table.find_all('td', limit=1)[0].find_parent('tr').find_all('td')
                
                # Add header row
                for header in headers:
                    markdown_table[0] += f" {header.get_text().strip()} |"
                
                # Add separator row
                markdown_table.append('|' + '---|' * len(headers))
                
                # Add data rows
                rows = table.find_all('tr')[1:] if headers else table.find_all('tr')
                for row in rows:
                    cells = row.find_all(['td', 'th'])
                    markdown_row = '|'
                    for cell in cells:
                        markdown_row += f" {cell.get_text().strip()} |"
                    markdown_table.append(markdown_row)
                
                # Replace table with markdown version
                table.replace_with(soup.new_string('\n'.join(markdown_table)))
            
            # Convert the cleaned HTML to markdown
            html_text = str(soup)
            markdown_content = markdown.markdown(html_text)
            
            # Clean up markdown
            # Remove extra newlines
            markdown_content = re.sub(r'\n{3,}', '\n\n', markdown_content)
            
            # Add title as heading if not present
            if not markdown_content.startswith('# '):
                markdown_content = f"# {title}\n\n{markdown_content}"
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert HTML file {file_path}: {str(e)}")
            raise ProcessingError(f"HTML conversion failed: {str(e)}") from e

class CSVConverter(BaseContentConverter):
    """Converts CSV content to markdown tables."""
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert CSV to markdown table.
        
        Args:
            file_path: Path to CSV file
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read CSV file
            df = pd.read_csv(file_path)
            
            # Extract metadata
            metadata = {
                'original_format': 'csv',
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': df.columns.tolist()
            }
            
            # Convert to markdown table
            markdown_lines = []
            
            # Add header
            header = '| ' + ' | '.join(str(col) for col in df.columns) + ' |'
            markdown_lines.append(header)
            
            # Add separator
            separator = '|' + '---|' * len(df.columns)
            markdown_lines.append(separator)
            
            # Add data rows
            for _, row in df.iterrows():
                markdown_row = '| ' + ' | '.join(str(cell) for cell in row) + ' |'
                markdown_lines.append(markdown_row)
            
            markdown_content = '\n'.join(markdown_lines)
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert CSV file {file_path}: {str(e)}")
            raise ProcessingError(f"CSV conversion failed: {str(e)}") from e

class JSONConverter(BaseContentConverter):
    """Converts JSON content to markdown code blocks."""
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert JSON to markdown code block.
        
        Args:
            file_path: Path to JSON file
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read JSON content
            with open(file_path, 'r', encoding='utf-8') as f:
                json_content = json.load(f)
            
            # Extract metadata
            metadata = {
                'original_format': 'json',
                'size': os.path.getsize(file_path),
                'top_level_keys': list(json_content.keys()) if isinstance(json_content, dict) else None,
                'array_length': len(json_content) if isinstance(json_content, list) else None
            }
            
            # Convert to pretty-printed JSON
            formatted_json = json.dumps(json_content, indent=2)
            
            # Wrap in markdown code block
            markdown_content = f"```json\n{formatted_json}\n```"
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert JSON file {file_path}: {str(e)}")
            raise ProcessingError(f"JSON conversion failed: {str(e)}") from e

class TextConverter(BaseContentConverter):
    """Converts text files to markdown code blocks."""

    def convert(self, content: Union[str, Path]) -> Tuple[str, Dict[str, Any]]:
        """Convert text to markdown code block.

        Args:
            content: Text content to convert or path to text file

        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read file if path provided
            if isinstance(content, Path):
                with open(content, 'r') as f:
                    text_content = f.read()
            else:
                text_content = content

            # Extract metadata
            metadata = {
                'original_format': 'text',
                'lines': text_content.count('\n') + 1,
                'characters': len(text_content)
            }

            # Check if content looks like YAML
            try:
                yaml.safe_load(text_content)
                is_yaml = True
            except:
                is_yaml = False

            # Wrap in appropriate markdown code block
            if is_yaml:
                markdown_content = f"```yaml\n{text_content}\n```"
                metadata['detected_format'] = 'yaml'
            else:
                markdown_content = f"```text\n{text_content}\n```"
                metadata['detected_format'] = 'plain_text'

            return markdown_content, metadata

        except Exception as e:
            self.logger.error(f"Failed to convert text content: {str(e)}")
            raise ProcessingError(f"Text conversion failed: {str(e)}") from e

class OfficeConverter(BaseContentConverter):
    """Converts Office documents to markdown."""
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert Office document to markdown.
        
        Args:
            file_path: Path to Office document
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        # Determine file type from extension
        ext = file_path.suffix.lower()
        
        if ext in ['.docx', '.doc']:
            return self._convert_word(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self._convert_excel(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._convert_powerpoint(file_path)
        else:
            raise ProcessingError(f"Unsupported Office format: {ext}")
            
    def _convert_word(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert Word document to markdown.
        
        Args:
            file_path: Path to Word document
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read document
            doc = docx.Document(file_path)
            
            # Extract metadata
            metadata = {
                'original_format': 'docx',
                'size': os.path.getsize(file_path),
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'sections': len(doc.sections)
            }
            
            # Convert content to markdown
            markdown_lines = []
            
            # Process paragraphs
            for paragraph in doc.paragraphs:
                if not paragraph.text.strip():
                    continue
                    
                # Handle heading styles
                style = paragraph.style.name.lower()
                if 'heading' in style:
                    level = int(style[-1]) if style[-1].isdigit() else 1
                    markdown_lines.append(f"{'#' * level} {paragraph.text}")
                else:
                    # Handle text formatting
                    text = paragraph.text
                    for run in paragraph.runs:
                        if run.bold:
                            text = text.replace(run.text, f"**{run.text}**")
                        if run.italic:
                            text = text.replace(run.text, f"*{run.text}*")
                    markdown_lines.append(text)
                
                markdown_lines.append('')  # Add blank line after paragraph
            
            # Process tables
            for table in doc.tables:
                # Table header
                header_row = []
                for cell in table.rows[0].cells:
                    header_row.append(cell.text.strip())
                markdown_lines.append('| ' + ' | '.join(header_row) + ' |')
                
                # Separator
                markdown_lines.append('|' + '---|' * len(header_row))
                
                # Table data
                for row in table.rows[1:]:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    markdown_lines.append('| ' + ' | '.join(row_data) + ' |')
                
                markdown_lines.append('')  # Add blank line after table
            
            markdown_content = '\n'.join(markdown_lines)
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert Word document {file_path}: {str(e)}")
            raise ProcessingError(f"Word conversion failed: {str(e)}") from e
            
    def _convert_excel(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert Excel document to markdown.
        
        Args:
            file_path: Path to Excel document
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read workbook
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            # Extract metadata
            metadata = {
                'original_format': 'xlsx',
                'size': os.path.getsize(file_path),
                'sheets': len(workbook.sheetnames),
                'active_sheet': workbook.active.title
            }
            
            markdown_lines = []
            
            # Process each worksheet
            for sheet in workbook.worksheets:
                # Add sheet title
                markdown_lines.append(f"## Sheet: {sheet.title}")
                markdown_lines.append("")
                
                # Get dimensions
                min_row = sheet.min_row
                max_row = sheet.max_row
                min_col = sheet.min_column
                max_col = sheet.max_column
                
                if min_row == 0 or min_col == 0:
                    continue  # Skip empty sheets
                
                # Get column headers (first row)
                headers = []
                for col in range(min_col, max_col + 1):
                    cell = sheet.cell(min_row, col)
                    headers.append(str(cell.value or ''))
                
                # Add table header
                markdown_lines.append('| ' + ' | '.join(headers) + ' |')
                markdown_lines.append('|' + '---|' * len(headers))
                
                # Add data rows
                for row in range(min_row + 1, max_row + 1):
                    row_data = []
                    for col in range(min_col, max_col + 1):
                        cell = sheet.cell(row, col)
                        value = cell.value
                        # Format numbers and dates appropriately
                        if isinstance(value, (int, float)):
                            value = f"{value:,}"
                        elif isinstance(value, datetime.datetime):
                            value = value.strftime('%Y-%m-%d %H:%M:%S')
                        elif isinstance(value, datetime.date):
                            value = value.strftime('%Y-%m-%d')
                        row_data.append(str(value or ''))
                    markdown_lines.append('| ' + ' | '.join(row_data) + ' |')
                
                markdown_lines.append("")  # Add blank line after table
            
            markdown_content = '\n'.join(markdown_lines)
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert Excel document {file_path}: {str(e)}")
            raise ProcessingError(f"Excel conversion failed: {str(e)}") from e
            
    def _convert_powerpoint(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert PowerPoint presentation to markdown.
        
        Args:
            file_path: Path to PowerPoint document
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            # Read presentation
            presentation = pptx.Presentation(file_path)
            
            # Extract metadata
            metadata = {
                'original_format': 'pptx',
                'size': os.path.getsize(file_path),
                'slides': len(presentation.slides)
            }
            
            markdown_lines = []
            
            # Process each slide
            for slide_number, slide in enumerate(presentation.slides, 1):
                # Add slide header
                markdown_lines.append(f"## Slide {slide_number}")
                markdown_lines.append("")
                
                # Get slide title
                if slide.shapes.title:
                    markdown_lines.append(f"### {slide.shapes.title.text}")
                    markdown_lines.append("")
                
                # Process shapes (text boxes, tables, etc.)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # Process text
                        for paragraph in shape.text_frame.paragraphs:
                            if not paragraph.text.strip():
                                continue
                            
                            # Handle text formatting
                            text = paragraph.text
                            for run in paragraph.runs:
                                if run.font.bold:
                                    text = text.replace(run.text, f"**{run.text}**")
                                if run.font.italic:
                                    text = text.replace(run.text, f"*{run.text}*")
                            markdown_lines.append(text)
                            markdown_lines.append("")
                    
                    elif shape.has_table:
                        # Process table
                        table = shape.table
                        
                        # Table header
                        header_row = []
                        for cell in table.rows[0].cells:
                            header_row.append(cell.text.strip())
                        markdown_lines.append('| ' + ' | '.join(header_row) + ' |')
                        
                        # Separator
                        markdown_lines.append('|' + '---|' * len(header_row))
                        
                        # Table data
                        for row in table.rows[1:]:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            markdown_lines.append('| ' + ' | '.join(row_data) + ' |')
                        
                        markdown_lines.append("")
                
                # Add slide notes if present
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    markdown_lines.append("#### Notes")
                    markdown_lines.append("")
                    markdown_lines.append(slide.notes_slide.notes_text_frame.text.strip())
                    markdown_lines.append("")
                
                markdown_lines.append("---")  # Add slide separator
                markdown_lines.append("")
            
            markdown_content = '\n'.join(markdown_lines)
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert PowerPoint document {file_path}: {str(e)}")
            raise ProcessingError(f"PowerPoint conversion failed: {str(e)}") from e

class PDFConverter(BaseContentConverter):
    """Converts PDF documents to markdown."""
    
    def __init__(self, temp_dir: Optional[Path] = None):
        """Initialize PDF converter.
        
        Args:
            temp_dir: Optional directory for temporary files. If not provided, uses system temp dir.
        """
        super().__init__()
        self.temp_dir = temp_dir or Path('/tmp')
    
    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert PDF document to markdown.
        
        Args:
            file_path: Path to PDF document
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            import fitz  # PyMuPDF
            
            # Open PDF document
            doc = fitz.open(file_path)
            
            # Extract metadata
            metadata = {
                'original_format': 'pdf',
                'size': os.path.getsize(file_path),
                'pages': len(doc),
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
                'keywords': doc.metadata.get('keywords', '')
            }
            
            content = []
            
            # Process each page
            for page_num, page in enumerate(doc, 1):
                content.append(f"\n## Page {page_num}\n")
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    content.append(text)
                
                # Extract images
                images = page.get_images()
                if images:
                    content.append("\n### Images\n")
                    for img_index, img in enumerate(images, 1):
                        xref = img[0]  # Cross-reference number
                        
                        # Get image data
                        base_image = doc.extract_image(xref)
                        if base_image:
                            # Get image info
                            image_data = base_image["image"]
                            image_ext = base_image["ext"]
                            
                            # Generate unique filename
                            image_filename = f"page{page_num}_image{img_index}.{image_ext}"
                            
                            # Save image to temp file
                            temp_path = self.temp_dir / image_filename
                            with open(temp_path, 'wb') as f:
                                f.write(image_data)
                            
                            # Convert image to markdown
                            image_converter = ImageConverter()
                            image_content, image_metadata = image_converter.convert(temp_path)
                            content.append(image_content)
                            
                            # Clean up temp file
                            temp_path.unlink()
                
                # Extract tables
                tables = page.find_tables()
                if tables:
                    content.append("\n### Tables\n")
                    for table in tables:
                        # Convert table to markdown format
                        rows = table.extract()
                        if rows:
                            # Add table header
                            content.append("| " + " | ".join(str(cell) for cell in rows[0]) + " |")
                            content.append("|" + "|".join("---" for _ in rows[0]) + "|")
                            
                            # Add table rows
                            for row in rows[1:]:
                                content.append("| " + " | ".join(str(cell) for cell in row) + " |")
                            content.append("")
            
            return "\n".join(content), metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert PDF document {file_path}: {str(e)}")
            raise ProcessingError(f"PDF conversion failed: {str(e)}") from e

class ImageConverter(BaseContentConverter):
    """Converts image files to markdown with descriptions."""
    
    def __init__(self, vision_api_key: Optional[str] = None):
        """Initialize image converter.

        Args:
            vision_api_key: Optional OpenAI Vision API key for image descriptions
        """
        super().__init__()
        self.vision_api_key = vision_api_key

    def convert(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Convert image to markdown with description.
        
        Args:
            file_path: Path to image file
            
        Returns:
            Tuple of (markdown content, metadata dict)
        """
        try:
            from PIL import Image
            import openai
            
            # Open image
            img = Image.open(file_path)
            
            # Extract metadata
            metadata = {
                'original_format': img.format.lower(),
                'size': os.path.getsize(file_path),
                'dimensions': f"{img.width}x{img.height}",
                'mode': img.mode,
                'dpi': img.info.get('dpi', 'Unknown')
            }
            
            # Default description
            description = "No AI description available"
            
            # Generate description using OpenAI Vision if available
            if self.vision_api_key:
                try:
                    openai.api_key = self.vision_api_key
                    
                    # Encode image to base64
                    import base64
                    from io import BytesIO
                    
                    # Convert to RGB if needed
                    if img.mode not in ('RGB', 'RGBA'):
                        img = img.convert('RGB')
                        
                    # Save to bytes
                    buffered = BytesIO()
                    img.save(buffered, format="JPEG")
                    img_str = base64.b64encode(buffered.getvalue()).decode()
                    
                    # Call OpenAI Vision API
                    response = openai.ChatCompletion.create(
                        model="gpt-4-vision-preview",
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": "Please describe this image in detail."},
                                    {
                                        "type": "image_url",
                                        "image_url": f"data:image/jpeg;base64,{img_str}"
                                    }
                                ]
                            }
                        ],
                        max_tokens=300
                    )
                    
                    description = response.choices[0].message.content
                    metadata['ai_description'] = description
                    
                except Exception as e:
                    self.logger.warning(f"Failed to generate AI description: {str(e)}")
            
            # Convert to markdown
            content = []
            
            # Add image reference
            content.append(f"![{file_path.name}]({file_path})")
            
            # Add description
            content.append(f"\n**Description:** {description}")
            
            # Add technical details
            content.append("\n**Technical Details:**")
            content.append(f"- Format: {metadata['original_format']}")
            content.append(f"- Dimensions: {metadata['dimensions']}")
            content.append(f"- Mode: {metadata['mode']}")
            content.append(f"- DPI: {metadata['dpi']}")
            content.append(f"- Size: {metadata['size']} bytes")
            
            return "\n".join(content), metadata
            
        except Exception as e:
            self.logger.error(f"Failed to convert image file {file_path}: {str(e)}")
            raise ProcessingError(f"Image conversion failed: {str(e)}") from e

class ContentConverterFactory:
    """Factory for content converters."""

    def __init__(self):
        """Initialize the factory."""
        self.logger = logging.getLogger(__name__)
        self._base_path = None
        self._anchor_ids = {}  # Changed from set to dict
        self._nav_links = {}
        self._references = {}

    def _validate_content(self, content: str) -> None:
        """Validate markdown content.
        
        Args:
            content: Content to validate
            
        Raises:
            ProcessingError: If content is invalid
        """
        if not content:
            raise ProcessingError("Content cannot be empty")
            
        # Check for headers
        if not re.search(r'^#', content, re.MULTILINE):
            self.logger.warning("Content has no headers")
            
        # Check table structure
        table_rows = re.findall(r'\|.*\|', content)
        if table_rows:
            col_counts = [len(row.split('|')) - 2 for row in table_rows]  # -2 for empty ends
            if len(set(col_counts)) > 1:
                raise ProcessingError("Table has inconsistent number of columns")
                
        # Check code blocks
        if content.count('```') % 2 != 0:
            raise ProcessingError("Unmatched code block markers")
            
        # Check image alt text
        if re.search(r'!\[\]', content):
            self.logger.warning("Image missing alt text")
            
    def _validate_metadata(self, metadata: Dict[str, Any]) -> None:
        """Validate metadata.
        
        Args:
            metadata: Metadata to validate
            
        Raises:
            ProcessingError: If metadata is invalid
        """
        if 'original_format' not in metadata:
            raise ProcessingError("Missing required metadata field: original_format")
            
        if 'size' in metadata and metadata['size'] <= 0:
            raise ProcessingError("File size must be positive")
            
        if 'dimensions' in metadata:
            if not re.match(r'^\d+x\d+$', str(metadata['dimensions'])):
                raise ProcessingError("Invalid dimensions format")
                
    def normalize_path(self, path: Union[str, Path]) -> str:
        """Normalize path relative to base path.
        
        Args:
            path: Path to normalize
            
        Returns:
            Normalized path string
        """
        try:
            path = Path(path)
            if not self._base_path:
                self.logger.warning("No base path set, returning absolute path")
                return str(path.resolve())
            
            # Convert to absolute path if relative
            if not path.is_absolute():
                path = (self._base_path / path).resolve()
            
            # Try to make relative to base path
            try:
                return str(path.relative_to(self._base_path))
            except ValueError:
                self.logger.warning(f"Path {path} is outside base path {self._base_path}")
                return str(path)
            
        except Exception as e:
            self.logger.warning(f"Failed to normalize path {path}: {str(e)}")
            return str(path)

    def generate_anchor_id(self, title: str) -> str:
        """Generate a unique anchor ID from a title."""
        # Convert to lowercase and replace special characters with hyphens
        anchor_id = title.lower()
        # Replace special characters with hyphens, preserving numbers
        anchor_id = re.sub(r'[^a-z0-9]+', '-', anchor_id)
        # Remove leading/trailing hyphens
        anchor_id = anchor_id.strip('-')
        # Ensure numbers are preserved correctly (e.g., 1.2.3 -> 123)
        anchor_id = re.sub(r'(\d+)\.(\d+)\.(\d+)', r'\1\2\3', anchor_id)
        
        # If this exact anchor ID already exists and it's for a different title
        if anchor_id in self._anchor_ids and self._anchor_ids[anchor_id].lower() != title.lower():
            base_id = anchor_id
            counter = 1
            while f"{base_id}-{counter}" in self._anchor_ids:
                counter += 1
            anchor_id = f"{base_id}-{counter}"
        
        # Store the anchor ID and its title
        self._anchor_ids[anchor_id] = title
        return anchor_id

    def resolve_relative_path(self, path: str, reference_path: Optional[Union[str, Path]] = None) -> Path:
        """Resolve a relative path against a base path."""
        # If it's a URL, return it as is
        if re.match(r'^[a-z]+://', path):
            return Path(path)
            
        try:
            # Convert to Path objects
            path_obj = Path(path)
            
            # If absolute path, verify it exists
            if path_obj.is_absolute():
                if not path_obj.exists():
                    raise ProcessingError(f"Path does not exist: {path_obj}")
                return path_obj
            
            # If reference path provided, resolve relative to it
            if reference_path:
                ref_path = Path(reference_path)
                if not ref_path.is_absolute() and self._base_path:
                    ref_path = self._base_path / ref_path
                resolved = (ref_path.parent / path_obj).resolve()
                if not resolved.exists():
                    raise ProcessingError(f"Path does not exist: {resolved}")
                return resolved
            
            # If no reference path but base path exists, resolve relative to base
            if self._base_path:
                resolved = (self._base_path / path_obj).resolve()
                if not resolved.exists():
                    raise ProcessingError(f"Path does not exist: {resolved}")
                return resolved
            
            # No base path set
            raise ProcessingError("No base path set")
            
        except Exception as e:
            if isinstance(e, ProcessingError):
                raise
            raise ProcessingError(f"Failed to resolve path '{path}': {str(e)}")

    def update_relative_paths(self, content: str, reference_path: Optional[Union[str, Path]] = None) -> str:
        """Update relative paths in content."""
        def replace_path(match):
            path = match.group(1)
            try:
                # Skip URLs
                if path.startswith(('http://', 'https://')):
                    return match.group(0)
                    
                resolved = self.resolve_relative_path(path, reference_path)
                if not resolved.exists():
                    self.logger.warning(f"Failed to resolve path {path}: File does not exist")
                    return match.group(0)
                return f'({self.normalize_path(resolved)})'
            except Exception as e:
                self.logger.warning(f"Failed to resolve path {path}: {str(e)}")
                return match.group(0)
                
        # Update image and link paths
        content = re.sub(r'\(([^)]+)\)', replace_path, content)
        return content

    def set_base_path(self, path: Union[str, Path]) -> None:
        """Set base path for resolving relative paths."""
        self._base_path = Path(path)

    def get_anchor_id(self, title: str) -> Optional[str]:
        """Get existing anchor ID for title."""
        for anchor_id, stored_title in self._anchor_ids.items():
            if stored_title.lower() == title.lower():
                return anchor_id
        return None

    def clear_anchor_ids(self) -> None:
        """Clear all anchor IDs."""
        self._anchor_ids.clear()

    def add_anchor_ids(self, content: str) -> str:
        """Add anchor IDs to all headers in the content."""
        # Reset anchor IDs for this content
        self._anchor_ids.clear()
        
        def replace_header(match):
            level = len(match.group(1))  # Number of # symbols
            title = match.group(2).strip()
            anchor_id = self.generate_anchor_id(title)
            return f"{match.group(1)} {title} {{#{anchor_id}}}"
        
        # Add anchor IDs to headers
        content = re.sub(r'^(#{1,6})\s+([^{#\n]+)(?:\s+{#[^}]+})?$',
                        replace_header,
                        content,
                        flags=re.MULTILINE)
        
        return content

    def update_navigation_links(self, content: str, file_path: str) -> str:
        """Update navigation links in content."""
        if not self._nav_links:
            return content
            
        nav_links = self._nav_links.get(file_path, {})
        if not nav_links:
            return content
            
        def replace_nav_link(match):
            link_type = match.group(1).lower()
            if link_type in nav_links:
                return f"<!-- {link_type}: {nav_links[link_type]} -->"
            return match.group(0)
            
        # Update navigation links
        content = re.sub(r'<!--\s*(prev|next|parent):\s*[^>]+\s*-->', replace_nav_link, content)
        return content

    def validate_navigation_links(self) -> List[str]:
        """Validate navigation links."""
        errors = []
        for file_path, links in self._nav_links.items():
            for link_type, target in links.items():
                try:
                    resolved = self.resolve_relative_path(target)
                    if not Path(resolved).exists():
                        errors.append(f"Navigation {link_type} file not found: {target}")
                except ProcessingError as e:
                    errors.append(f"Failed to resolve {link_type} link {target}: {str(e)}")
                    
                # Check for broken prev/next chain
                if link_type == 'next':
                    next_file = target
                    next_links = self._nav_links.get(next_file, {})
                    if 'prev' not in next_links or next_links['prev'] != file_path:
                        errors.append(f"Broken next/prev chain between {file_path} and {next_file}")
                        
        return errors

    def get_converter(self, file_path: Path) -> BaseContentConverter:
        """Get appropriate converter for file type.
        
        Args:
            file_path: Path to file to convert
            
        Returns:
            Appropriate converter instance
            
        Raises:
            ProcessingError: If no converter available for file type
        """
        suffix = file_path.suffix.lower()
        
        converters = {
            '.html': HTMLConverter,
            '.csv': CSVConverter,
            '.json': JSONConverter,
            '.txt': TextConverter,
            '.docx': OfficeConverter,
            '.pdf': PDFConverter,
            '.jpg': ImageConverter,
            '.jpeg': ImageConverter,
            '.png': ImageConverter,
            '.gif': ImageConverter
        }
        
        if suffix not in converters:
            raise ProcessingError(f"No converter available for file type: {suffix}")
            
        return converters[suffix]()
        
    def wrap_with_markers(self, content: str, file_path: Path) -> str:
        """Wrap content with attachment markers.
        
        Args:
            content: Content to wrap
            file_path: Path to file (for filename in marker)
            
        Returns:
            Content wrapped with attachment markers
        """
        filename = file_path.name
        return f"--==ATTACHMENT_BLOCK: {filename}==--\n{content}\n--==ATTACHMENT_BLOCK_END==--"
        
    def detect_attachment_blocks(self, content: str) -> List[Dict[str, Any]]:
        """Detect attachment blocks in content.
        
        Args:
            content: Content to search for attachment blocks
            
        Returns:
            List of attachment block dictionaries
        """
        blocks = []
        lines = content.split('\n')
        current_block = None
        
        for i, line in enumerate(lines, 1):
            if match := re.match(r'--==ATTACHMENT_BLOCK: (.+?)==--', line):
                if current_block:
                    self.logger.warning("Unclosed attachment block")
                current_block = {
                    'filename': match.group(1),
                    'content': [],
                    'start_line': i,
                    'start_pos': content.find(line)
                }
            elif line.strip() == '--==ATTACHMENT_BLOCK_END==--' and current_block:
                current_block['end_line'] = i
                current_block['end_pos'] = content.find(line) + len(line)
                current_block['content'] = '\n'.join(current_block['content'])
                blocks.append(current_block)
                current_block = None
            elif current_block:
                current_block['content'].append(line)
                
        if current_block:
            self.logger.warning("Unclosed attachment block")
            
        return blocks
        
    def extract_attachment_metadata(self, block: Dict[str, Any]) -> Dict[str, Any]:
        """Extract metadata from attachment block.
        
        Args:
            block: Attachment block dictionary
            
        Returns:
            Metadata dictionary
        """
        metadata = {
            'filename': block['filename'],
            'file_type': Path(block['filename']).suffix[1:],
            'line_range': (block['start_line'], block['end_line']),
            'byte_range': (block['start_pos'], block['end_pos']),
            'content_size': len(block['content']),
            'has_references': bool(re.search(r'\[\[.+?\]\]', block['content']))
        }
        return metadata
        
    def validate_attachment_block(self, block: Dict[str, Any]) -> None:
        """Validate attachment block.
        
        Args:
            block: Attachment block dictionary
            
        Raises:
            ProcessingError: If block is invalid
        """
        required_fields = ['filename', 'content', 'start_line', 'end_line', 'start_pos', 'end_pos']
        for field in required_fields:
            if field not in block:
                raise ProcessingError(f"Missing required field: {field}")
                
        if not block['filename']:
            raise ProcessingError("Empty filename")
            
        if block['start_line'] > block['end_line']:
            raise ProcessingError("Invalid line range")
            
        if not block['content']:
            raise ProcessingError("Empty content")
            
    def update_cross_references(self, content: str) -> str:
        """Update cross-references in content.
        
        Args:
            content: Content to update
            
        Returns:
            Content with updated cross-references
        """
        def update_ref(match):
            title = match.group(1)
            anchor_id = self.get_anchor_id(title)
            if not anchor_id:
                self.logger.warning(f"Cross-reference target not found: {title}")
                anchor_id = "missing-ref"
            return f'[{title}](#{anchor_id})'
            
        return re.sub(r'\[\[(.+?)\]\]', update_ref, content)
        
    def clear_navigation_links(self) -> None:
        """Clear all navigation links."""
        self._nav_links.clear()
        
    def add_navigation_links(self, content: str, file_path: str) -> str:
        """Add navigation links to content.
        
        Args:
            content: Content to add links to
            file_path: Path of current file
            
        Returns:
            Content with navigation links added
        """
        if file_path not in self._nav_links:
            return content
            
        links = self._nav_links[file_path]
        nav_section = "\n\n## Navigation\n"
        
        if 'prev' in links:
            nav_section += f"← [Previous]({links['prev']}) "
        if 'next' in links:
            nav_section += f"→ [Next]({links['next']}) "
        if 'parent' in links:
            nav_section += f"↑ [Up]({links['parent']})"
            
        return content + nav_section
        
    def extract_navigation_links(self, content: str, file_path: str) -> Dict[str, str]:
        """Extract navigation links from content.
        
        Args:
            content: Content to extract links from
            file_path: Path of current file
            
        Returns:
            Dictionary of navigation links
        """
        links = {}
        patterns = {
            'prev': r'<!-- prev: (.+?) -->',
            'next': r'<!-- next: (.+?) -->',
            'parent': r'<!-- parent: (.+?) -->'
        }
        
        for link_type, pattern in patterns.items():
            if match := re.search(pattern, content):
                links[link_type] = match.group(1)
                
        return links

# Now we need to create a factory to handle different file types