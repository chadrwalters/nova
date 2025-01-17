"""Document converter class."""
import logging
from pathlib import Path
from typing import Any, cast

import html2text
import pdfplumber
import pypandoc
import yaml
from openpyxl import load_workbook
from pptx import Presentation

from .datamodel import Document, InputFormat

logger = logging.getLogger(__name__)


class DocumentConverter:
    """Converts documents between formats."""

    def __init__(self) -> None:
        """Initialize the converter."""
        self._html2text = html2text.HTML2Text()
        self._html2text.ignore_links = False  # Keep links for reference
        self._html2text.ignore_images = False  # Keep image references
        self._html2text.ignore_tables = False
        self._html2text.body_width = 0
        self._html2text.single_line_break = True
        self._html2text.protect_links = True
        self._html2text.unicode_snob = True  # Use Unicode characters
        self._html2text.wrap_links = False  # Don't wrap links
        self._html2text.inline_links = True  # Use inline links
        self._html2text.emphasis_mark = "*"  # Use * for emphasis
        self._html2text.strong_mark = "**"  # Use ** for strong
        self._html2text.google_doc = False  # Not a Google Doc
        self._html2text.ul_item_mark = "-"  # Use - for unordered lists
        self._html2text.default_image_alt = ""  # Empty alt text for images
        self._html2text.skip_internal_links = False  # Keep internal links
        self._html2text.pad_tables = True  # Add padding to tables
        self._html2text.escape_snob = True  # Don't escape special characters
        self._html2text.images_to_alt = True  # Use alt text for images
        self._html2text.images_with_size = False  # Don't include image sizes
        self._html2text.decode_errors = "strict"  # Strict decoding

    def _parse_frontmatter(self, content: str) -> tuple[dict[str, Any], str]:
        """Parse YAML frontmatter from content.

        Args:
            content: Content to parse.

        Returns:
            Tuple of metadata dict and remaining content.
        """
        if not content.startswith("---\n"):
            return {}, content

        parts = content.split("---\n", 2)
        if len(parts) < 3:
            return {}, content

        try:
            metadata = yaml.safe_load(parts[1])
            if not isinstance(metadata, dict):
                metadata = {}
            return metadata, parts[2]
        except yaml.YAMLError:
            return {}, content

    def convert(self, doc: Document) -> Document:
        """Convert a document to markdown format."""
        if doc.format in (InputFormat.MD, InputFormat.TEXT):
            # Parse frontmatter for markdown files
            if doc.format == InputFormat.MD:
                metadata, content = self._parse_frontmatter(doc.content)
                return Document(
                    content=content,
                    format=InputFormat.MD,
                    metadata=metadata,
                    title=metadata.get("title"),
                    tags=metadata.get("tags", []),
                )
            # Text files are already in a markdown-compatible format
            return Document(content=doc.content, format=InputFormat.MD)

        content = doc.content
        fmt = doc.format

        if fmt == InputFormat.HTML:
            logger.info("Converting HTML content: %s", content)
            content = cast(str, pypandoc.convert_text(content, "markdown", format="html"))
            logger.info("Converted to markdown: %s", content)
        elif fmt in (InputFormat.RST, InputFormat.ASCIIDOC, InputFormat.LATEX):
            content = cast(str, pypandoc.convert_text(content, "markdown", format=str(fmt)))
        else:
            raise ValueError(f"Unsupported conversion from {fmt} to markdown")

        return Document(content=content, format=InputFormat.MD)

    def convert_file(self, file_path: Path, fmt: InputFormat) -> Document:
        """Convert a file to markdown format."""
        if fmt in (InputFormat.MD, InputFormat.TEXT):
            with open(file_path, encoding="utf-8") as f:
                content = f.read()
            if fmt == InputFormat.MD:
                metadata, content = self._parse_frontmatter(content)
                return Document(
                    content=content,
                    format=InputFormat.MD,
                    metadata=metadata,
                    title=metadata.get("title"),
                    tags=metadata.get("tags", []),
                    source_path=file_path,
                )
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt == InputFormat.HTML:
            logger.info("Converting HTML file: %s", file_path)
            with open(file_path, encoding="utf-8") as f:
                content = f.read()
            logger.info("HTML content: %s", content)
            content = cast(str, pypandoc.convert_text(content, "markdown", format="html"))
            logger.info("Converted to markdown: %s", content)
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt == InputFormat.PDF:
            with pdfplumber.open(file_path) as pdf:
                pages = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                content = "\n\n".join(pages)
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt == InputFormat.DOCX:
            content = cast(str, pypandoc.convert_file(str(file_path), "markdown"))
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt == InputFormat.XLSX:
            wb = load_workbook(filename=file_path)
            sheets = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    rows.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
                if rows:
                    sheets.append(f"## {sheet.title}\n\n" + "\n".join(rows))
            content = "\n\n".join(sheets)
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt == InputFormat.PPTX:
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = getattr(shape, "text", "")
                        if text:
                            texts.append(text)
                if texts:
                    slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            content = "\n\n".join(slides)
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        if fmt in (InputFormat.RST, InputFormat.ASCIIDOC, InputFormat.LATEX):
            with open(file_path, encoding="utf-8") as f:
                content = f.read()
            content = cast(str, pypandoc.convert_text(content, "markdown", format=str(fmt)))
            return Document(content=content, format=InputFormat.MD, source_path=file_path)

        raise ValueError(f"Unsupported conversion from {fmt} to markdown")
