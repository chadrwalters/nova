"""Markdown handling components for Nova processors."""

from pathlib import Path
from typing import Dict, Any, Optional, List
import re
from datetime import datetime
from urllib.parse import unquote, quote
import os
import tempfile
import subprocess
import aiofiles

from . import MarkdownComponent
from ...core.config import NovaConfig, ProcessorConfig
from ...core.errors import ProcessingError, ImageProcessingError
from ...core.logging import get_logger
from ...core.image_processor import ImageProcessor

class MarkitdownHandler(MarkdownComponent):
    """Handles markdown processing using markitdown."""
    
    def __init__(self, processor_config: ProcessorConfig, nova_config: NovaConfig, image_processor: Optional[ImageProcessor] = None):
        """Initialize handler.
        
        Args:
            processor_config: Processor-specific configuration
            nova_config: Global Nova configuration
            image_processor: Optional image processor instance
        """
        super().__init__(processor_config, nova_config)
        self.logger = get_logger(self.__class__.__name__)
        
        # Initialize or use provided image processor
        self.image_processor = image_processor or ImageProcessor(config=processor_config.get('image', {}))
        
        # Add component-specific stats
        self.stats.update({
            'files_processed': 0,
            'images_processed': 0,
            'links_updated': 0
        })
    
    async def process_markdown(self, content: str, source_path: Path) -> str:
        """Process markdown content, handling images and other elements."""
        self.logger.info(f"Processing markdown file: {source_path}")
        self.logger.info("Processing images in markdown content")
        
        # Process line by line to maintain original formatting
        lines = content.split('\n')
        processed_lines = []
        
        # Regex pattern for markdown image syntax
        image_pattern = r'!\[(.*?)\]\((.*?)\)'
        
        for line in lines:
            self.logger.info(f"Processing line: {line}")
            if re.search(image_pattern, line):
                self.logger.info("Line matched pattern")
                # Process images in the line
                processed_line = await self._process_line_images(line, image_pattern, source_path)
                processed_lines.append(processed_line)
            else:
                self.logger.info("Line did not match pattern")
                processed_lines.append(line)
        
        return '\n'.join(processed_lines)
    
    async def _process_line_images(self, line: str, pattern: str, source_path: Path) -> str:
        """Process all images in a line of text.
        
        Args:
            line: Line of text to process
            pattern: Regex pattern for finding images
            source_path: Path to source markdown file
            
        Returns:
            Processed line with updated image references
        """
        result = line
        for match in re.finditer(pattern, line):
            replacement = await self._process_image_match(match, source_path)
            result = result.replace(match.group(0), replacement)
        return result
    
    async def _process_image_match(self, match: re.Match, source_path: Path) -> str:
        """Process a matched image in markdown content.
        
        Args:
            match: Regex match object
            source_path: Path to source markdown file
            
        Returns:
            Updated markdown image reference
        """
        try:
            alt_text = match.group(1)
            image_path_str = unquote(match.group(2))
            image_path = Path(image_path_str)
            
            # Make path absolute if relative
            if not image_path.is_absolute():
                image_path = source_path.parent / image_path
            
            # Skip if image doesn't exist
            if not image_path.exists():
                self.logger.warning(f"Image not found: {image_path}")
                return match.group(0)  # Return original unchanged
            
            # Skip if not a supported format
            if image_path.suffix.lower() not in self.image_processor.supported_formats:
                self.logger.warning(f"Unsupported image format: {image_path}")
                return match.group(0)  # Return original unchanged
            
            # Process image
            try:
                # Convert HEIC if needed
                if image_path.suffix.lower() in ['.heic', '.heif']:
                    self.logger.info(f"Converting HEIC image: {image_path}")
                    image_path = await self.image_processor.convert_heic_to_jpg(image_path)
                
                # Optimize if configured
                if self.processor_config.get('optimize_images', False):
                    preset = self.processor_config.get('optimization_preset', 'web')
                    image_path = await self.image_processor.optimize_image(image_path, preset)
                
                # Extract metadata
                metadata = await self.image_processor.extract_metadata(image_path)
                
                # Update alt text with dimensions if not already present
                if not re.search(r'\(\d+x\d+\)', alt_text):
                    dimensions = metadata['basic']['dimensions']
                    if not alt_text:
                        alt_text = image_path.stem
                    alt_text = f"{alt_text} ({dimensions[0]}x{dimensions[1]})"
                
                # Update stats
                self.stats['images_processed'] += 1
                
                # Return updated markdown reference
                return f"![{alt_text}]({quote(str(image_path.relative_to(source_path.parent)))})"
                
            except Exception as e:
                self.logger.error(f"Error processing image {image_path}: {str(e)}")
                return match.group(0)  # Return original unchanged
            
        except Exception as e:
            self.logger.error(f"Error in _process_image_match: {str(e)}")
            return match.group(0)  # Return original unchanged
    
    async def convert_to_markdown(self, input_path: Path, output_path: Optional[Path] = None) -> str:
        """Convert a file to markdown format.
        
        Args:
            input_path: Path to input file
            output_path: Optional path to save markdown output
            
        Returns:
            Converted markdown content
        """
        try:
            self.logger.info(f"Converting file to markdown: {input_path}")
            
            # Get file type from suffix
            suffix = input_path.suffix.lower()
            
            # Convert based on file type
            if suffix in ['.docx', '.doc']:
                self.logger.info(f"Converting Word document: {input_path}")
                content = await self._convert_docx(input_path, output_path)
            elif suffix in ['.pptx', '.ppt']:
                self.logger.info(f"Converting PowerPoint presentation: {input_path}")
                content = await self._convert_pptx(input_path, output_path)
            elif suffix in ['.xlsx', '.xls']:
                self.logger.info(f"Converting Excel workbook: {input_path}")
                content = await self._convert_xlsx(input_path, output_path)
            elif suffix == '.pdf':
                self.logger.info(f"Converting PDF document: {input_path}")
                content = await self._convert_pdf(input_path, output_path)
            elif suffix == '.csv':
                self.logger.info(f"Converting CSV file: {input_path}")
                content = await self._convert_csv(input_path, output_path)
            elif suffix in ['.txt', '.json', '.env', '.html']:
                self.logger.info(f"Converting text file: {input_path}")
                # For text files, read directly
                async with aiofiles.open(input_path, 'r', encoding='utf-8') as f:
                    content = await f.read()
            elif suffix in self.image_processor.supported_formats:
                self.logger.info(f"Creating markdown reference for image: {input_path}")
                # For images, create a markdown image reference with metadata
                metadata = await self.image_processor.extract_metadata(input_path)
                dimensions = metadata['basic']['dimensions']
                content = f"![{input_path.stem} ({dimensions[0]}x{dimensions[1]})]({input_path.name})"
            else:
                error = f"Unsupported file format '{suffix}' for file: {input_path}"
                self.logger.error(error)
                raise ProcessingError(error)
            
            # Save to output file if specified
            if output_path:
                self.logger.info(f"Writing markdown output to: {output_path}")
                output_path.parent.mkdir(parents=True, exist_ok=True)
                async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                    await f.write(content)
            
            self.logger.info(f"Successfully converted {input_path} to markdown")
            return content
            
        except Exception as e:
            error = f"Failed to convert {input_path} (format: {suffix}): {str(e)}"
            self.logger.error(error)
            return f"# Error Converting {input_path.name}\n\nFile type: {suffix}\nError: {str(e)}\n"
    
    async def _convert_docx(self, input_path: Path, output_path: Path) -> str:
        """Convert DOCX to markdown."""
        from docx import Document
        
        self.logger.info(f"Opening Word document: {input_path}")
        doc = Document(str(input_path))
        content = []
        
        self.logger.info(f"Extracting paragraphs from: {input_path}")
        for para in doc.paragraphs:
            content.append(para.text)
        
        self.logger.info(f"Successfully extracted {len(content)} paragraphs from: {input_path}")
        return '\n\n'.join(content)
    
    async def _convert_pptx(self, input_path: Path, output_path: Path) -> str:
        """Convert PPTX to markdown."""
        from pptx import Presentation
        
        self.logger.info(f"Opening PowerPoint presentation: {input_path}")
        prs = Presentation(str(input_path))
        content = []
        
        self.logger.info(f"Processing slides from: {input_path}")
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            self.logger.debug(f"Processing slide {i} from: {input_path}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            content.append('\n'.join(slide_content))
        
        self.logger.info(f"Successfully processed {len(prs.slides)} slides from: {input_path}")
        return '\n\n'.join(content)
    
    async def _convert_xlsx(self, input_path: Path, output_path: Path) -> str:
        """Convert XLSX to markdown."""
        from openpyxl import load_workbook
        from datetime import datetime
        
        def format_cell_value(value):
            """Format cell value for markdown."""
            if value is None:
                return ''
            elif isinstance(value, datetime):
                return value.strftime('%Y-%m-%d')  # Format dates as YYYY-MM-DD
            else:
                return str(value).strip()
        
        self.logger.info(f"Opening Excel workbook: {input_path}")
        wb = load_workbook(str(input_path))
        content = []
        
        self.logger.info(f"Processing {len(wb.sheetnames)} sheets from: {input_path}")
        for sheet in wb:
            self.logger.debug(f"Processing sheet '{sheet.title}' from: {input_path}")
            # Add sheet name as header
            content.append(f"# {sheet.title}\n")
            
            # Process each sheet
            row_idx = 1
            tables_found = 0
            while row_idx <= sheet.max_row:
                # Find the start of a table (non-empty row)
                if any(cell.value for cell in sheet[row_idx]):
                    tables_found += 1
                    self.logger.debug(f"Found table {tables_found} in sheet '{sheet.title}' at row {row_idx}")
                    table_content = []
                    header_cells = list(sheet[row_idx])
                    
                    # Look ahead for merged header rows
                    next_row_idx = row_idx + 1
                    while next_row_idx <= sheet.max_row:
                        next_row = list(sheet[next_row_idx])
                        if not any(cell.value for cell in next_row):
                            break
                        if all(isinstance(cell.value, (str, type(None))) for cell in next_row):
                            header_cells.extend(next_row)
                            next_row_idx += 1
                        else:
                            break
                    
                    # Build header from merged rows
                    headers = []
                    for cell in header_cells:
                        value = format_cell_value(cell.value)
                        if value:
                            headers.append(value)
                    
                    # Skip if no valid headers
                    if not headers:
                        row_idx = next_row_idx
                        continue
                    
                    # Add header row
                    table_content.append(' | '.join(headers))
                    
                    # Add separator row
                    table_content.append('|'.join(['---' for _ in headers]))
                    
                    # Process data rows
                    data_start_row = next_row_idx
                    rows_processed = 0
                    while data_start_row <= sheet.max_row:
                        row = list(sheet[data_start_row])
                        row_content = []
                        is_empty_row = True
                        
                        # Process each cell in the row
                        for cell in row[:len(headers)]:  # Only process cells up to header count
                            value = format_cell_value(cell.value)
                            row_content.append(value)
                            if value:
                                is_empty_row = False
                        
                        # Skip empty rows
                        if is_empty_row:
                            data_start_row += 1
                            continue
                        
                        # Pad row if needed
                        while len(row_content) < len(headers):
                            row_content.append('')
                        
                        # Add data row
                        table_content.append(' | '.join(row_content))
                        data_start_row += 1
                        rows_processed += 1
                    
                    self.logger.debug(f"Processed {rows_processed} rows in table {tables_found}")
                    # Add table to content with spacing
                    content.extend(['', *table_content, ''])
                    row_idx = data_start_row
                else:
                    row_idx += 1
            
            self.logger.debug(f"Found {tables_found} tables in sheet '{sheet.title}'")
            # Add blank line between sheets
            content.append('')
        
        self.logger.info(f"Successfully processed workbook with {len(wb.sheetnames)} sheets from: {input_path}")
        return '\n'.join(content)
    
    async def _convert_pdf(self, input_path: Path, output_path: Path) -> str:
        """Convert PDF to markdown."""
        import warnings
        from pypdf import PdfReader
        
        self.logger.info(f"Opening PDF document: {input_path}")
        # Suppress PyPDF deprecation warning about ARC4
        with warnings.catch_warnings():
            warnings.filterwarnings(
                'ignore',
                message='ARC4 has been moved to cryptography.hazmat.decrepit.*',
                category=DeprecationWarning
            )
            reader = PdfReader(str(input_path))
            content = []
            
            self.logger.info(f"Processing {len(reader.pages)} pages from: {input_path}")
            for i, page in enumerate(reader.pages, 1):
                self.logger.debug(f"Extracting text from page {i} of: {input_path}")
                content.append(page.extract_text())
            
            self.logger.info(f"Successfully extracted text from {len(reader.pages)} pages of: {input_path}")
            return '\n\n'.join(content)
    
    async def _convert_csv(self, input_path: Path, output_path: Path) -> str:
        """Convert CSV to markdown."""
        import csv
        import chardet
        
        # Detect encoding
        self.logger.info(f"Detecting encoding for CSV file: {input_path}")
        with open(input_path, 'rb') as f:
            raw = f.read()
            result = chardet.detect(raw)
            encoding = result['encoding']
        self.logger.debug(f"Detected encoding {encoding} for: {input_path}")
        
        # Read CSV with detected encoding
        content = []
        self.logger.info(f"Reading CSV file with encoding {encoding}: {input_path}")
        with open(input_path, 'r', encoding=encoding) as f:
            reader = csv.reader(f)
            rows_processed = 0
            for row in reader:
                content.append(' | '.join(row))
                rows_processed += 1
        
        self.logger.info(f"Successfully processed {rows_processed} rows from: {input_path}")
        return '\n'.join(content)
    
    async def _validate_image_file(self, image_path: Path) -> tuple[bool, str]:
        """Validate an image file.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        self.logger.debug(f"Validating image file: {image_path}")
        
        if not image_path.exists():
            self.logger.debug(f"Image file not found: {image_path}")
            return False, f"Image file not found: {image_path}"
            
        if not image_path.is_file():
            self.logger.debug(f"Path exists but is not a file: {image_path}")
            return False, f"Path exists but is not a file: {image_path}"
            
        # Check if it's a supported image format
        suffix = image_path.suffix.lower()
        if suffix not in self.image_processor.supported_formats:
            self.logger.debug(f"Unsupported image format '{suffix}' for file: {image_path}")
            return False, f"Unsupported image format '{suffix}' for file: {image_path}"
            
        # Check if file is readable
        try:
            image_path.open('rb').close()
            self.logger.debug(f"Successfully validated image file: {image_path}")
            return True, ""
        except Exception as e:
            self.logger.debug(f"File is not readable: {image_path} ({str(e)})")
            return False, f"File is not readable: {image_path} ({str(e)})"

class ConsolidationHandler(MarkdownComponent):
    """Handles markdown consolidation."""
    
    def __init__(self, processor_config: ProcessorConfig, nova_config: NovaConfig):
        """Initialize handler.
        
        Args:
            processor_config: Processor-specific configuration
            nova_config: Global Nova configuration
        """
        super().__init__(processor_config, nova_config)
        self.logger = get_logger(self.__class__.__name__)
        
        # Add component-specific stats
        self.stats.update({
            'files_consolidated': 0,
            'total_files': 0,
            'total_size': 0
        })
    
    def process_markdown(self, content: str, source_path: Path) -> str:
        """Process markdown content.
        
        Args:
            content: Markdown content to process
            source_path: Path to source file
            
        Returns:
            Processed markdown content
        """
        # This handler doesn't modify individual markdown files
        return content
    
    def consolidate_markdown(self, input_files: List[Path], output_path: Path) -> Path:
        """Consolidate markdown files.
        
        Args:
            input_files: List of input files to consolidate
            output_path: Path to output file
            
        Returns:
            Path to consolidated file
        """
        try:
            # Sort files by date in filename (YYYYMMDD format)
            date_pattern = r'(\d{8})'
            sorted_files = sorted(
                input_files,
                key=lambda x: re.search(date_pattern, x.name).group(1) if re.search(date_pattern, x.name) else '00000000'
            )
            
            # Read and process each file
            consolidated_content = []
            for file_path in sorted_files:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    
                    # Add file header
                    header = f"\n\n## {file_path.stem}\n\n"
                    consolidated_content.append(header)
                    
                    # Add processed content
                    consolidated_content.append(content)
                    
                    # Update stats
                    self.stats['files_consolidated'] += 1
                    self.stats['total_size'] += len(content)
            
            # Write consolidated content
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(consolidated_content))
            
            # Update total files stat
            self.stats['total_files'] = len(input_files)
            
            return output_path
            
        except Exception as e:
            raise ProcessingError(f"Failed to consolidate markdown: {e}") from e