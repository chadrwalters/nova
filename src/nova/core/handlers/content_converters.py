"""Content converters for different file formats."""

import csv
import io
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, Type, Optional
import tempfile
import os

import aiofiles
import docx
import openpyxl
import markitdown
from pptx import Presentation
import yaml

class BaseContentConverter(ABC):
    """Base class for content converters."""
    
    @abstractmethod
    async def convert(self, file_path: Path) -> str:
        """Convert file content to markdown.
        
        Args:
            file_path: Path to file to convert
            
        Returns:
            Markdown representation of file content
            
        Raises:
            Exception: If conversion fails
        """
        pass

class MarkdownConverter(BaseContentConverter):
    """Converter for markdown files."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert markdown file.
        
        Just reads the file as is since it's already markdown.
        Validates the content and extracts any metadata.
        
        Args:
            file_path: Path to markdown file
            
        Returns:
            Markdown content
        """
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
            
        # Check for YAML frontmatter
        if content.startswith('---'):
            try:
                # Find end of frontmatter
                end = content.find('---', 3)
                if end != -1:
                    # Parse frontmatter
                    frontmatter = yaml.safe_load(content[3:end])
                    # Reconstruct content with validated frontmatter
                    content = f"---\n{yaml.dump(frontmatter)}---\n{content[end+4:]}"
            except yaml.YAMLError:
                pass  # Invalid frontmatter, return as is
                
        return content

class DocxConverter(BaseContentConverter):
    """Converter for Word documents."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert Word document to markdown.
        
        Args:
            file_path: Path to Word document
            
        Returns:
            Markdown representation of document
        """
        doc = docx.Document(file_path)
        
        # Convert paragraphs to markdown
        markdown = []
        for para in doc.paragraphs:
            # Skip empty paragraphs
            if not para.text.strip():
                continue
                
            # Convert paragraph style to markdown
            if para.style.name.startswith('Heading'):
                level = para.style.name[-1]
                markdown.append(f"{'#' * int(level)} {para.text}")
            else:
                markdown.append(para.text)
            
            # Add newline after paragraph
            markdown.append('')
        
        return '\n'.join(markdown)

class PptxConverter(BaseContentConverter):
    """Converter for PowerPoint presentations."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert PowerPoint presentation to markdown.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Markdown representation of presentation
        """
        prs = Presentation(file_path)
        
        # Convert slides to markdown
        markdown = []
        for i, slide in enumerate(prs.slides, 1):
            # Add slide header
            markdown.append(f"# Slide {i}")
            markdown.append('')
            
            # Convert shapes to text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown.append(shape.text)
                    markdown.append('')
            
            # Add slide separator
            markdown.append('---')
            markdown.append('')
        
        return '\n'.join(markdown)

class XlsxConverter(BaseContentConverter):
    """Converter for Excel spreadsheets."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert Excel spreadsheet to markdown.
        
        Args:
            file_path: Path to Excel file
            
        Returns:
            Markdown representation of spreadsheet
        """
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet = wb.active
        
        # Convert worksheet to markdown table
        markdown = []
        
        # Get dimensions
        max_row = sheet.max_row
        max_col = sheet.max_column
        
        # Create header row
        header = []
        for col in range(1, max_col + 1):
            cell = sheet.cell(row=1, column=col)
            header.append(str(cell.value or ''))
        markdown.append('| ' + ' | '.join(header) + ' |')
        
        # Add separator row
        markdown.append('|' + '|'.join(['---'] * max_col) + '|')
        
        # Add data rows
        for row in range(2, max_row + 1):
            data = []
            for col in range(1, max_col + 1):
                cell = sheet.cell(row=row, column=col)
                data.append(str(cell.value or ''))
            markdown.append('| ' + ' | '.join(data) + ' |')
        
        return '\n'.join(markdown)

class PdfConverter(BaseContentConverter):
    """Converter for PDF documents."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert PDF document to markdown.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Markdown representation of document
        """
        # Convert PDF to markdown using markitdown
        converter = markitdown.MarkItDown()
        result = converter.convert_local(str(file_path), output_format="markdown")
        return result.text_content

class CsvConverter(BaseContentConverter):
    """Converter for CSV files."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert CSV file to markdown.
        
        Args:
            file_path: Path to CSV file
            
        Returns:
            Markdown representation of CSV data
        """
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
        content = None
        
        for encoding in encodings:
            try:
                async with aiofiles.open(file_path, 'r', encoding=encoding) as f:
                    content = await f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise ValueError(f"Could not decode {file_path} with any of the attempted encodings")
        
        # Parse CSV content
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        
        if not rows:
            return ''
        
        # Convert to markdown table
        markdown = []
        
        # Add header row with escaped pipes
        header = [cell.replace('|', '\\|') for cell in rows[0]]
        markdown.append('| ' + ' | '.join(header) + ' |')
        
        # Add separator row
        markdown.append('|' + '|'.join(['---'] * len(rows[0])) + '|')
        
        # Add data rows with escaped pipes and special character handling
        for row in rows[1:]:
            escaped_row = [cell.replace('|', '\\|').replace('\n', '<br>') for cell in row]
            markdown.append('| ' + ' | '.join(escaped_row) + ' |')
        
        return '\n'.join(markdown)

class HtmlConverter(BaseContentConverter):
    """Converter for HTML files."""
    
    async def convert(self, file_path: Path) -> str:
        """Convert HTML file to markdown.
        
        Args:
            file_path: Path to HTML file
            
        Returns:
            Markdown representation of HTML content
        """
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
            
        # Create a temporary file for the HTML content
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', encoding='utf-8', delete=False) as temp_file:
            temp_file.write(content)
            temp_path = temp_file.name
            
        try:
            # Convert HTML to markdown using markitdown
            converter = markitdown.MarkItDown()
            result = converter.convert_local(temp_path, output_format="markdown")
            return result.text_content
        finally:
            # Clean up temporary file
            os.unlink(temp_path)

class ContentConverterFactory:
    """Factory for creating content converters."""
    
    _converters: Dict[str, Type[BaseContentConverter]] = {
        'md': MarkdownConverter,
        'docx': DocxConverter,
        'doc': DocxConverter,
        'pptx': PptxConverter,
        'ppt': PptxConverter,
        'xlsx': XlsxConverter,
        'xls': XlsxConverter,
        'pdf': PdfConverter,
        'csv': CsvConverter,
        'html': HtmlConverter,
        'htm': HtmlConverter
    }
    
    @classmethod
    def get_converter(cls, file_type: str) -> Optional[BaseContentConverter]:
        """Get converter for file type.
        
        Args:
            file_type: File extension without dot
            
        Returns:
            Converter instance if supported, None otherwise
        """
        converter_class = cls._converters.get(file_type.lower())
        if converter_class:
            return converter_class()
        return None