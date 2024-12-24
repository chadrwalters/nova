"""Office document processor module for Nova document processor."""

import os
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
import json
from docx import Document
import openpyxl
from pptx import Presentation
import fitz  # PyMuPDF

from ..core.base import BaseProcessor
from ..core.config import ProcessorConfig, OfficeConfig
from ..core.paths import PathsConfig
from ..core.errors import ProcessingError
from ..core.logging import get_logger

class OfficeProcessor(BaseProcessor):
    """Processes office documents."""

    def __init__(self, config: OfficeConfig, paths: PathsConfig):
        """Initialize office processor.
        
        Args:
            config: Office processor configuration
            paths: Path configuration
        """
        super().__init__(config, paths)
        self.cache = {}
        self._load_cache()

    def process(self, input_path: Path) -> Dict[str, Any]:
        """Process office document.
        
        Args:
            input_path: Path to office document
            
        Returns:
            Dictionary containing processing results
            
        Raises:
            ProcessingError: If processing fails
        """
        try:
            self._ensure_output_dir()
            output_path = self._get_output_path(input_path)
            
            # Check cache
            if self._is_cached(input_path):
                self._log_progress(f"Using cached version of {input_path}")
                return self._get_cached_result(input_path)
            
            # Process document based on type
            suffix = input_path.suffix.lower()
            if suffix in ['.docx', '.doc']:
                content, metadata = self._process_word(input_path)
            elif suffix in ['.xlsx', '.xls']:
                content, metadata = self._process_excel(input_path)
            elif suffix in ['.pptx', '.ppt']:
                content, metadata = self._process_powerpoint(input_path)
            elif suffix == '.pdf':
                content, metadata = self._process_pdf(input_path)
            else:
                raise ProcessingError(f"Unsupported file type: {suffix}")
            
            # Save processed content
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path.with_suffix('.md'), 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Save metadata
            metadata_path = self._get_metadata_path(input_path)
            self._save_metadata(metadata_path, metadata)
            
            result = {
                'input_path': str(input_path),
                'output_path': str(output_path),
                'metadata_path': str(metadata_path),
                'metadata': metadata,
                'status': 'success'
            }
            
            # Update cache
            self._cache_result(input_path, result)
            
            return result
            
        except Exception as e:
            context = {'input_path': str(input_path)}
            self._handle_error(e, context)

    def _process_word(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process Word document.
        
        Args:
            file_path: Path to Word document
            
        Returns:
            Tuple of (markdown content, metadata)
        """
        doc = Document(file_path)
        content = []
        
        # Extract metadata
        metadata = {
            'format': 'docx',
            'paragraphs': len(doc.paragraphs),
            'sections': len(doc.sections),
            'tables': len(doc.tables)
        }
        
        # Process paragraphs
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
                
            # Handle heading styles
            style = para.style.name.lower()
            if 'heading' in style:
                level = int(style[-1]) if style[-1].isdigit() else 1
                content.append(f"{'#' * level} {para.text}")
            else:
                content.append(para.text)
            content.append("")
        
        # Process tables
        for table in doc.tables:
            # Table header
            header_row = []
            for cell in table.rows[0].cells:
                header_row.append(cell.text.strip())
            content.append('| ' + ' | '.join(header_row) + ' |')
            
            # Separator
            content.append('|' + '---|' * len(header_row))
            
            # Table data
            for row in table.rows[1:]:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                content.append('| ' + ' | '.join(row_data) + ' |')
            content.append("")
        
        return "\n".join(content), metadata

    def _process_excel(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process Excel document.
        
        Args:
            file_path: Path to Excel document
            
        Returns:
            Tuple of (markdown content, metadata)
        """
        wb = openpyxl.load_workbook(file_path, data_only=True)
        content = []
        
        # Extract metadata
        metadata = {
            'format': 'xlsx',
            'sheets': len(wb.sheetnames),
            'active_sheet': wb.active.title
        }
        
        # Process each worksheet
        for sheet in wb.worksheets:
            content.append(f"## Sheet: {sheet.title}")
            content.append("")
            
            # Get dimensions
            min_row = sheet.min_row
            max_row = sheet.max_row
            min_col = sheet.min_column
            max_col = sheet.max_column
            
            if min_row == 0 or min_col == 0:
                continue
            
            # Get column headers
            headers = []
            for col in range(min_col, max_col + 1):
                cell = sheet.cell(min_row, col)
                headers.append(str(cell.value or ''))
            
            # Add table header
            content.append('| ' + ' | '.join(headers) + ' |')
            content.append('|' + '---|' * len(headers))
            
            # Add data rows
            for row in range(min_row + 1, max_row + 1):
                row_data = []
                for col in range(min_col, max_col + 1):
                    cell = sheet.cell(row, col)
                    row_data.append(str(cell.value or ''))
                content.append('| ' + ' | '.join(row_data) + ' |')
            content.append("")
        
        return "\n".join(content), metadata

    def _process_powerpoint(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process PowerPoint document.
        
        Args:
            file_path: Path to PowerPoint document
            
        Returns:
            Tuple of (markdown content, metadata)
        """
        prs = Presentation(file_path)
        content = []
        
        # Extract metadata
        metadata = {
            'format': 'pptx',
            'slides': len(prs.slides)
        }
        
        # Process each slide
        for idx, slide in enumerate(prs.slides, 1):
            content.append(f"## Slide {idx}")
            content.append("")
            
            # Get slide title
            if slide.shapes.title:
                content.append(f"### {slide.shapes.title.text}")
                content.append("")
            
            # Process shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            content.append(para.text)
                            content.append("")
                
                elif shape.has_table:
                    table = shape.table
                    # Table header
                    header_row = []
                    for cell in table.rows[0].cells:
                        header_row.append(cell.text.strip())
                    content.append('| ' + ' | '.join(header_row) + ' |')
                    
                    # Separator
                    content.append('|' + '---|' * len(header_row))
                    
                    # Table data
                    for row in table.rows[1:]:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        content.append('| ' + ' | '.join(row_data) + ' |')
                    content.append("")
            
            content.append("---")
            content.append("")
        
        return "\n".join(content), metadata

    def _process_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process PDF document.
        
        Args:
            file_path: Path to PDF document
            
        Returns:
            Tuple of (markdown content, metadata)
        """
        doc = fitz.open(file_path)
        content = []
        
        # Extract metadata
        metadata = {
            'format': 'pdf',
            'pages': len(doc),
            'title': doc.metadata.get('title', ''),
            'author': doc.metadata.get('author', ''),
            'subject': doc.metadata.get('subject', '')
        }
        
        # Process each page
        for page_num, page in enumerate(doc, 1):
            content.append(f"## Page {page_num}")
            content.append("")
            
            # Extract text
            text = page.get_text()
            if text.strip():
                content.append(text)
                content.append("")
            
            # Extract tables
            tables = page.find_tables()
            if tables:
                content.append("### Tables")
                content.append("")
                for table in tables:
                    rows = table.extract()
                    if rows:
                        # Table header
                        content.append('| ' + ' | '.join(str(cell) for cell in rows[0]) + ' |')
                        content.append('|' + '---|' * len(rows[0]))
                        
                        # Table data
                        for row in rows[1:]:
                            content.append('| ' + ' | '.join(str(cell) for cell in row) + ' |')
                        content.append("")
        
        return "\n".join(content), metadata

    def _get_metadata_path(self, input_path: Path) -> Path:
        """Get path for metadata file.
        
        Args:
            input_path: Input document path
            
        Returns:
            Path for metadata file
        """
        rel_path = input_path.relative_to(self.paths.input_dir)
        return Path(self.paths.office_dirs['metadata']) / f"{rel_path.stem}.json"

    def _save_metadata(self, metadata_path: Path, metadata: Dict[str, Any]) -> None:
        """Save metadata to file.
        
        Args:
            metadata_path: Path to save metadata
            metadata: Metadata to save
        """
        metadata_path.parent.mkdir(parents=True, exist_ok=True)
        with open(metadata_path, 'w') as f:
            json.dump(metadata, f, indent=2, default=str)

    def _is_cached(self, input_path: Path) -> bool:
        """Check if document is cached.
        
        Args:
            input_path: Input document path
            
        Returns:
            True if document is cached
        """
        if str(input_path) not in self.cache:
            return False
            
        cached = self.cache[str(input_path)]
        if not os.path.exists(cached['output_path']):
            return False
            
        input_mtime = os.path.getmtime(input_path)
        cache_mtime = os.path.getmtime(cached['output_path'])
        
        return input_mtime <= cache_mtime

    def _get_cached_result(self, input_path: Path) -> Dict[str, Any]:
        """Get cached processing result.
        
        Args:
            input_path: Input document path
            
        Returns:
            Cached result dictionary
        """
        return self.cache[str(input_path)]

    def _cache_result(self, input_path: Path, result: Dict[str, Any]) -> None:
        """Cache processing result.
        
        Args:
            input_path: Input document path
            result: Result to cache
        """
        self.cache[str(input_path)] = result
        self._save_cache()

    def _load_cache(self) -> None:
        """Load cache from file."""
        cache_path = self.paths.office_dirs['cache'] / 'office_cache.json'
        if cache_path.exists():
            try:
                with open(cache_path, 'r') as f:
                    self.cache = json.load(f)
            except Exception as e:
                self.logger.warning(f"Failed to load cache: {e}")
                self.cache = {}

    def _save_cache(self) -> None:
        """Save cache to file."""
        cache_path = self.paths.office_dirs['cache'] / 'office_cache.json'
        cache_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            with open(cache_path, 'w') as f:
                json.dump(self.cache, f, indent=2, default=str)
        except Exception as e:
            self.logger.warning(f"Failed to save cache: {e}")