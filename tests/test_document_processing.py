"""Test document processing functionality."""

import pytest
import tempfile
from pathlib import Path
from reportlab.pdfgen import canvas
from docx import Document
from docx.shared import Inches
import pypandoc
from pptx import Presentation

def create_test_pdf(output_path: Path) -> None:
    """Create a test PDF with known content."""
    c = canvas.Canvas(str(output_path))
    c.drawString(100, 750, "Test PDF Document")
    c.drawString(100, 700, "This is a test paragraph")
    c.drawString(100, 650, "With multiple lines")
    c.drawString(100, 600, "To verify extraction")
    c.save()

def create_test_docx(output_path: Path) -> None:
    """Create a test DOCX with formatted content."""
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a normal paragraph.')
    doc.add_paragraph('Bullet point 1', style='List Bullet')
    doc.add_paragraph('Bullet point 2', style='List Bullet')
    
    table = doc.add_table(rows=2, cols=2)
    table.style = 'Table Grid'
    cells = table.rows[0].cells
    cells[0].text = 'Header 1'
    cells[1].text = 'Header 2'
    cells = table.rows[1].cells
    cells[0].text = 'Data 1'
    cells[1].text = 'Data 2'
    
    doc.save(str(output_path))

def create_test_rtf(output_path: Path) -> None:
    """Create a test RTF file."""
    content = """# Test RTF Document
    
    This is a test paragraph.
    
    * Bullet point 1
    * Bullet point 2
    
    **Bold text** and *italic text*.
    """
    pypandoc.convert_text(content, 'rtf', format='markdown', outputfile=str(output_path))

def create_test_odt(output_path: Path) -> None:
    """Create a test ODT file."""
    content = """# Test ODT Document
    
    This is a test paragraph.
    
    1. Numbered item 1
    2. Numbered item 2
    
    | Header 1 | Header 2 |
    |----------|----------|
    | Data 1   | Data 2   |
    """
    pypandoc.convert_text(content, 'odt', format='markdown', outputfile=str(output_path))

def create_test_pptx(output_path: Path) -> None:
    """Create a test PowerPoint file."""
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "For Testing Purposes"
    
    # Content slide
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Test Slide'
    tf = body_shape.text_frame
    tf.text = 'First bullet point'
    p = tf.add_paragraph()
    p.text = 'Second bullet point'
    p.level = 1
    
    prs.save(str(output_path))

@pytest.fixture
def test_files():
    """Create test files for each format."""
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # Create test files
        pdf_path = temp_path / "test.pdf"
        create_test_pdf(pdf_path)
        
        docx_path = temp_path / "test.docx"
        create_test_docx(docx_path)
        
        rtf_path = temp_path / "test.rtf"
        create_test_rtf(rtf_path)
        
        odt_path = temp_path / "test.odt"
        create_test_odt(odt_path)
        
        pptx_path = temp_path / "test.pptx"
        create_test_pptx(pptx_path)
        
        yield {
            'pdf': pdf_path,
            'docx': docx_path,
            'rtf': rtf_path,
            'odt': odt_path,
            'pptx': pptx_path
        }

def test_pdf_extraction(test_files):
    """Test PDF text extraction quality."""
    pdf_path = test_files['pdf']
    
    # Read the PDF and verify content
    with open(pdf_path, 'rb') as f:
        from PyPDF2 import PdfReader
        reader = PdfReader(f)
        text = reader.pages[0].extract_text()
        
        assert "Test PDF Document" in text
        assert "This is a test paragraph" in text
        assert "With multiple lines" in text
        assert "To verify extraction" in text

def test_docx_formatting(test_files):
    """Test DOCX formatting preservation."""
    docx_path = test_files['docx']
    
    # Read the DOCX and verify content/formatting
    doc = Document(docx_path)
    
    # Check heading (can be Title or Heading style)
    assert doc.paragraphs[0].text == 'Test Document'
    assert doc.paragraphs[0].style.name in ['Title', 'Heading 1']
    
    # Check bullet points
    bullet_points = [p for p in doc.paragraphs if p.style.name == 'List Bullet']
    assert len(bullet_points) == 2
    assert bullet_points[0].text == 'Bullet point 1'
    assert bullet_points[1].text == 'Bullet point 2'
    
    # Check table
    assert len(doc.tables) == 1
    table = doc.tables[0]
    assert table.rows[0].cells[0].text == 'Header 1'
    assert table.rows[1].cells[1].text == 'Data 2'

def test_rtf_odt_support(test_files):
    """Test RTF and ODT support."""
    # Convert RTF to markdown and verify
    rtf_md = pypandoc.convert_file(str(test_files['rtf']), 'markdown')
    assert 'Test RTF Document' in rtf_md
    assert 'Bullet point 1' in rtf_md
    assert 'Bullet point 2' in rtf_md
    
    # Convert ODT to markdown and verify
    odt_md = pypandoc.convert_file(str(test_files['odt']), 'markdown')
    assert 'Test ODT Document' in odt_md
    assert 'Numbered item 1' in odt_md
    assert 'Numbered item 2' in odt_md
    assert 'Header 1' in odt_md and 'Data 1' in odt_md

def test_pptx_handling(test_files):
    """Test PowerPoint handling."""
    prs = Presentation(test_files['pptx'])
    
    # Check slides count
    assert len(prs.slides) == 2
    
    # Check title slide
    title_slide = prs.slides[0]
    assert title_slide.shapes.title.text == "Test Presentation"
    
    # Check content slide
    content_slide = prs.slides[1]
    assert content_slide.shapes.title.text == "Test Slide"
    
    # Check bullet points
    text_frame = content_slide.shapes.placeholders[1].text_frame
    assert "First bullet point" in text_frame.text
    assert "Second bullet point" in text_frame.text 