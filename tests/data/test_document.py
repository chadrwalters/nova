"""Test document generation utilities."""

from pathlib import Path
from reportlab.pdfgen import canvas
from docx import Document
from docx.shared import Inches
from pptx import Presentation
import pypandoc


def create_test_pdf(output_path: Path) -> None:
    """Create a test PDF file."""
    c = canvas.Canvas(str(output_path))
    
    # Page 1
    c.drawString(100, 750, "Test PDF Document")
    c.drawString(100, 700, "This is page 1")
    c.drawString(100, 650, "Some test content")
    
    # Page 2
    c.showPage()
    c.drawString(100, 750, "Page 2")
    c.drawString(100, 700, "More test content")
    
    c.save()


def create_test_docx(output_path: Path) -> None:
    """Create a test DOCX file."""
    doc = Document()
    
    doc.add_heading('Test Document', 0)
    
    # First section
    doc.add_heading('First Section', level=1)
    doc.add_paragraph('This is the first section of the test document.')
    
    # Bullet points
    doc.add_paragraph('First bullet point', style='List Bullet')
    doc.add_paragraph('Second bullet point', style='List Bullet')
    
    # Second section
    doc.add_heading('Second Section', level=1)
    doc.add_paragraph('This is the second section with a table.')
    
    # Add table
    table = doc.add_table(rows=2, cols=2)
    table.style = 'Table Grid'
    cells = table.rows[0].cells
    cells[0].text = 'Header 1'
    cells[1].text = 'Header 2'
    cells = table.rows[1].cells
    cells[0].text = 'Data 1'
    cells[1].text = 'Data 2'
    
    doc.save(str(output_path))


def create_test_rtf(output_path: Path) -> None:
    """Create a test RTF file."""
    markdown_content = """
# Test RTF Document

## Section 1
This is the first section of the test document.

## Section 2
This is the second section with some formatting:
* Bullet point 1
* Bullet point 2

**Bold text** and *italic text*.
"""
    pypandoc.convert_text(
        markdown_content,
        'rtf',
        format='markdown',
        outputfile=str(output_path)
    )


def create_test_odt(output_path: Path) -> None:
    """Create a test ODT file."""
    markdown_content = """
# Test ODT Document

## First Section
This is the first section of the test document.

## Second Section
This is the second section with a list:
1. First item
2. Second item

And a table:

| Column 1 | Column 2 |
|----------|----------|
| Data 1   | Data 2   |
"""
    pypandoc.convert_text(
        markdown_content,
        'odt',
        format='markdown',
        outputfile=str(output_path)
    )


def create_test_pptx(output_path: Path) -> None:
    """Create a test PPTX file."""
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created for testing"
    
    # Bullet points slide
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Bullet Points'
    
    tf = body_shape.text_frame
    tf.text = 'First bullet point'
    p = tf.add_paragraph()
    p.text = 'Second bullet point'
    p.level = 1
    
    # Table slide
    table_slide = prs.slides.add_slide(prs.slide_layouts[2])
    shapes = table_slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Table Example'
    
    rows = 2
    cols = 2
    left = Inches(2.0)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)
    
    # Fill table
    table.cell(0, 0).text = 'Header 1'
    table.cell(0, 1).text = 'Header 2'
    table.cell(1, 0).text = 'Data 1'
    table.cell(1, 1).text = 'Data 2'
    
    prs.save(str(output_path)) 